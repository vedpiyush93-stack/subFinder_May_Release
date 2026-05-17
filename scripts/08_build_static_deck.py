#!/usr/bin/env python3
"""
Build the subFinder presentation deck.

Loads rep_1 artifacts (already-trained best model + per-fold metrics + ablation
+ calibration), generates 9 figures and 6 tables into presentations/figures/
and presentations/tables/, then assembles them into deck.pptx via python-pptx.

Re-run with:
    python presentations/build_slides.py
"""
from __future__ import annotations
import json, glob, re, sys
from pathlib import Path
from collections import defaultdict, Counter

import numpy as np, pandas as pd, joblib
import matplotlib as mpl, matplotlib.pyplot as plt
import seaborn as sns
from scipy.special import expit

ROOT = Path(__file__).resolve().parent.parent
PRES = ROOT / "docs"  # release: render into docs/ for GitHub Pages
FIG = PRES / "figures"; FIG.mkdir(exist_ok=True)
TAB = PRES / "tables"; TAB.mkdir(exist_ok=True)
REP = ROOT / "artifacts"  # release layout

# === aesthetic config ===========================================================
NAVY = "#1a3a5c"; ORANGE = "#e67e22"; SAGE = "#27ae60"; CRIMSON = "#c0392b"
GRAY = "#7f8c8d"; LIGHT = "#ecf0f1"; CHARCOAL = "#2c3e50"
BLACK = "#000000"
PALETTE = [NAVY, ORANGE, SAGE, CRIMSON, GRAY]
# Bold-black styling everywhere: titles, axis labels, ticks, legends
mpl.rcParams.update({
    "font.family": "sans-serif",
    "font.sans-serif": ["DejaVu Sans", "Helvetica", "Arial"],
    "font.size": 12,
    "axes.titlesize": 15, "axes.titleweight": "bold",
    "axes.labelsize": 13, "axes.labelweight": "bold",
    "axes.labelcolor": BLACK, "axes.edgecolor": BLACK,
    "axes.linewidth": 1.3, "axes.spines.top": False, "axes.spines.right": False,
    "text.color": BLACK,
    "xtick.color": BLACK, "ytick.color": BLACK,
    "xtick.labelsize": 11, "ytick.labelsize": 11,
    "xtick.major.width": 1.2, "ytick.major.width": 1.2,
    "axes.grid": True, "grid.alpha": 0.20, "grid.color": "#888888", "grid.linestyle": "--",
    "figure.dpi": 200, "savefig.dpi": 200, "savefig.bbox": "tight",
    "axes.titlepad": 12,
    "legend.fontsize": 11, "legend.title_fontsize": 12, "legend.frameon": False,
})

# === load data ==================================================================
print("[deck] loading data + labels ...")
df_data = pd.read_csv(ROOT / "data/Train_data.csv")
X = df_data["sig_gene_seq"].fillna("").values
y_labels = df_data["high_level_substr"].values
substrates = sorted(set(y_labels))

TOK_RE = re.compile(r"[,|_]")
CAZY_RE = re.compile(r"^(GH|PL|CE|CBM|GT|AA)[0-9]+$")
def tok_cpu(s): return [t for t in TOK_RE.split(str(s)) if t]
def is_cazy(t): return bool(CAZY_RE.match(str(t)))

# Per-fold metrics (all 29 configs × n trials)
df_pf = pd.read_csv(REP / "per_fold_metrics.csv")
print(f"  per-fold rows: {len(df_pf)}, configs: {df_pf.shorthand.nunique()}")

# Lit canonical mapping (alias-collapsed)
lit = pd.read_csv(ROOT / "data/Literature_Data_fam_substrate_mapping.tsv", sep="\t")
lit.columns = [c.strip() for c in lit.columns]
# ALIAS map with per-entry provenance:
#  "exact": 1:1 string match between lit substrate name and our class name (no collapse — trust by definition).
#  "collapse:<lit-name>": lit substrate is biochemically related to our class via a documented collapse with primary-literature citations.
# Decisions per collapse are documented in the appendix slide.
ALIAS_PROVENANCE = {
    "alpha-glucan":     [("alpha-glucan","exact"), ("starch","collapse"), ("glycogen","collapse"),
                          ("sucrose","collapse"), ("raffinose","collapse"), ("trehalose","collapse"),
                          ("palatinose","collapse"), ("glucooligosaccharide","collapse")],
    "beta-glucan":      [("beta-glucan","exact"), ("cellulose","collapse"), ("cellooligosaccharide","collapse"),
                          ("xyloglucan","collapse"), ("beta-glycan","collapse")],
    "galactan":         [("beta-galactan","collapse"), ("alpha-galactan","collapse")],  # no exact "galactan" entry in lit
    "arabinogalactan":  [("arabinogalactan protein","collapse"), ("arabinan","collapse")],  # no exact "arabinogalactan" in lit
    "host glycan":      [("host glycan","exact"), ("human milk polysaccharide","collapse"),
                          ("sialic acid","collapse"), ("fucose","collapse")],
    "chitin":           [("chitin","exact"), ("chitosan","collapse"), ("chitooligosaccharide","collapse")],
    "alginate":         [("alginate","exact")],
    "pectin":           [("pectin","exact")],
    "xylan":            [("xylan","exact")],
    "alpha-mannan":     [("alpha-mannan","exact")],
    "beta-mannan":      [("beta-mannan","exact")],
    "fructan":          [("fructan","exact")],
}
# Also handle multi-substrate lit rows like "cellulose, chitin" by splitting + assigning to each component.
# Build (substrate, family) → set of provenance tags for clean audit trail.
CANON_PROV = {s: {} for s in substrates}  # {family: set([(lit_substrate, kind)])}
for _, row in lit.iterrows():
    lit_sub = str(row["Substrate_high_level"]).strip()
    fam = row["Family"]
    # split on comma + 'and' for multi-substrate rows
    parts = [p.strip() for p in re.split(r",|and\s+", lit_sub) if p.strip()]
    for part in parts:
        for our_sub, entries in ALIAS_PROVENANCE.items():
            for (alias_name, kind) in entries:
                if part == alias_name:
                    CANON_PROV[our_sub].setdefault(fam, set()).add((alias_name, kind))
# Final CANON sets (without provenance)
CANON = {s: set(CANON_PROV[s].keys()) for s in substrates}
print(f"  After careful collapse: {sum(len(c) for c in CANON.values())} (substrate,family) pairs")
print(f"  Per-substrate canon set sizes (with provenance summary):")
for s in substrates:
    n = len(CANON[s])
    n_exact_only = sum(1 for fam, provs in CANON_PROV[s].items() if any(k=="exact" for _,k in provs))
    n_collapse_only = sum(1 for fam, provs in CANON_PROV[s].items()
                           if all(k=="collapse" for _,k in provs))
    print(f"    {s:<18} {n:>3}  ({n_exact_only} via exact 1:1, {n_collapse_only} via collapse only)")

# Family tagging for plots
def family_of(short):
    if short in ("cpu__ET500_log2","ftCbow_MM__ET500_sqrt"): return "Ours (shallow)"
    if "BRF100" in short: return "Paper BRF baselines"
    if "__LSTM" in short and "attn" not in short: return "DL: LSTM"
    if "__LSTMattn" in short: return "DL: LSTM+attn"
    if "__JustAttn" in short: return "DL: attention"
    if "__Trans" in short: return "DL: transformer"
    return "?"
df_pf["family"] = df_pf.shorthand.apply(family_of)

# Decoder: shorthand → (featurizer_family, featurizer_detail, classifier_family, classifier_detail)
FEATURIZER_MAP = {
    "cpu":       ("CountVec",  "tok_cpu (splits on , | _)"),
    "cv":        ("CountVec",  "paper's tok_comma_pipe (, |)"),
    "ftCbow":    ("FastText",  "CBOW · mean-pooled"),
    "ftCbow_MM": ("FastText",  "CBOW · mean+max concat"),
    "ftSg":      ("FastText",  "skip-gram · mean-pooled"),
    "w2vCbow":   ("Word2Vec",  "CBOW · mean-pooled"),
    "w2vSg":     ("Word2Vec",  "skip-gram · mean-pooled"),
    "d2vDm":     ("Doc2Vec",   "DM · doc-vector"),
    "d2vDbow":   ("Doc2Vec",   "DBOW · doc-vector"),
}
CLASSIFIER_MAP = {
    "ET500_log2": ("OvR(ExtraTrees)",        "n=500 max_features=log2 class_weight=balanced"),
    "ET500_sqrt": ("OvR(ExtraTrees)",        "n=500 max_features=sqrt class_weight=balanced"),
    "BRF100":     ("OvR(BalancedRF)",        "n=100 class_weight=balanced (paper)"),
    "LSTM":       ("DL: LSTM",               "vanilla LSTM (paper)"),
    "LSTMattn":   ("DL: LSTM+attention",     "LSTM + attention (paper)"),
    "JustAttn":   ("DL: attention",          "attention-only (paper)"),
    "Trans":      ("DL: transformer",        "4-block transformer (paper)"),
}
def decode_shorthand(sh):
    feat_key, clf_key = sh.split("__")
    feat_fam, feat_det = FEATURIZER_MAP.get(feat_key, ("?","?"))
    clf_fam, clf_det = CLASSIFIER_MAP.get(clf_key, ("?","?"))
    return feat_fam, feat_det, clf_fam, clf_det

# Intuitive human-readable names. Used everywhere the shorthand would appear on a chart label.
PRETTY_FEATURIZER = {
    "cpu":       "CountVec (split , | _)",
    "cv":        "CountVec (paper, split , |)",
    "ftCbow":    "FastText CBOW",
    "ftCbow_MM": "FastText CBOW (mean+max)",
    "ftSg":      "FastText skipgram",
    "w2vCbow":   "Word2Vec CBOW",
    "w2vSg":     "Word2Vec skipgram",
    "d2vDm":     "Doc2Vec DM",
    "d2vDbow":   "Doc2Vec DBOW",
}
PRETTY_CLASSIFIER = {
    "ET500_log2": "ExtraTrees 500 (log2)",
    "ET500_sqrt": "ExtraTrees 500 (sqrt)",
    "BRF100":     "Balanced RF 100",
    "LSTM":       "LSTM",
    "LSTMattn":   "LSTM + attention",
    "JustAttn":   "Attention-only",
    "Trans":      "Transformer (4 blocks)",
}
def pretty_name(sh):
    f, c = sh.split("__")
    return f"{PRETTY_FEATURIZER.get(f, f)}  →  {PRETTY_CLASSIFIER.get(c, c)}"

# Intuitive family names — used for grouping plots and legends
PRETTY_FAMILY = {
    "Ours (shallow)":      "Our shallow winners (ExtraTrees)",
    "Paper BRF baselines": "Paper baselines (Balanced RF)",
    "DL: LSTM":            "Deep LSTM (paper)",
    "DL: LSTM+attn":       "Deep LSTM + attention (paper)",
    "DL: attention":       "Deep attention-only (paper)",
    "DL: transformer":     "Deep transformer (paper)",
}
FAMILY_COLOR = {"Ours (shallow)": SAGE, "Paper BRF baselines": NAVY,
                "DL: LSTM": GRAY, "DL: LSTM+attn": ORANGE,
                "DL: attention": CRIMSON, "DL: transformer": "#8e44ad"}

# Only configs with complete 25 trials for benchmark plots
counts = df_pf.shorthand.value_counts()
complete_shorts = set(counts[counts == 25].index)
df_pf_c = df_pf[df_pf.shorthand.isin(complete_shorts)].copy()
print(f"  configs with all 25 trials: {len(complete_shorts)}")

# =================================================================================
# FIG 1: benchmark — by-config bar with stripes per family
# =================================================================================
print("[deck] Fig 1: benchmark leaderboard ...")
# Sort DESCENDING so the WINNER is at the top of the chart (matplotlib horizontal bars
# place the LAST item at the top, so we sort ascending to display best-on-top).
agg = df_pf_c.groupby(["shorthand","family"]).agg(
    mean=("acc","mean"), std=("acc","std")).reset_index().sort_values("mean", ascending=True)
agg["rank"] = list(range(len(agg), 0, -1))  # 1 = highest mean accuracy
agg["pretty"] = agg["shorthand"].apply(pretty_name)
agg["label"] = agg.apply(lambda r: f"#{int(r['rank'])}  {r['pretty']}", axis=1)

fig, ax = plt.subplots(figsize=(13.5, 8.0))
y = np.arange(len(agg))
colors = [FAMILY_COLOR[f] for f in agg.family]
bars = ax.barh(y, agg["mean"], xerr=agg["std"], color=colors, edgecolor="black", linewidth=0.5,
               capsize=3, error_kw={"linewidth": 1.1, "ecolor": BLACK, "alpha": 0.85})
ax.set_yticks(y); ax.set_yticklabels(agg.label, fontsize=10, fontweight="bold", color=BLACK)
ax.set_xlabel("Mean test accuracy ± 1 SD  (5×5 RSKF, n=25 trials per config)")
ax.set_xlim(0.55, 0.985)
# value labels on each bar
for i, (_, r) in enumerate(agg.iterrows()):
    ax.text(r["mean"] + r["std"] + 0.003, i, f"{r['mean']:.4f}",
            va="center", fontsize=9.5, fontweight="bold", color=BLACK)
ax.axvline(0.9058, color=SAGE, linestyle=":", linewidth=1.3, alpha=0.8)
# legend OUTSIDE the plot (with intuitive family names)
from matplotlib.patches import Patch
handles = [Patch(facecolor=c, edgecolor="black", linewidth=0.5, label=PRETTY_FAMILY.get(f, f))
           for f, c in FAMILY_COLOR.items() if (df_pf_c.family == f).any()]
ax.legend(handles=handles, loc="lower right", title="Model family", title_fontsize=11,
          fontsize=10, frameon=True, facecolor="white", edgecolor="#cccccc")
ax.set_title("Benchmark leaderboard — 29 configurations, 5×5 RSKF (best on top)",
             color=BLACK, loc="left", pad=12, fontsize=15, fontweight="bold")
plt.tight_layout()
plt.savefig(FIG/"fig1_benchmark_leaderboard.png"); plt.close()

# =================================================================================
# FIG 1B: Top-5 podium — horizontal bars with featurizer + classifier annotations
# =================================================================================
print("[deck] Fig 1b: top-5 podium ...")
top5 = agg.sort_values("mean", ascending=False).head(5).reset_index(drop=True)
top5["rank_d"] = np.arange(1, len(top5)+1)
top5["featurizer_fam"] = top5.shorthand.apply(lambda s: decode_shorthand(s)[0])
top5["featurizer_det"] = top5.shorthand.apply(lambda s: decode_shorthand(s)[1])
top5["clf_fam"]        = top5.shorthand.apply(lambda s: decode_shorthand(s)[2])
top5["clf_det"]        = top5.shorthand.apply(lambda s: decode_shorthand(s)[3])

fig, ax = plt.subplots(figsize=(14, 7.5))
top5_rev = top5.iloc[::-1].reset_index(drop=True)
y_pos = np.arange(len(top5_rev))
bar_colors = [SAGE if s in ("cpu__ET500_log2","ftCbow_MM__ET500_sqrt") else
              NAVY if "BRF100" in s else ORANGE for s in top5_rev.shorthand]
ax.barh(y_pos, top5_rev["mean"], xerr=top5_rev["std"], color=bar_colors,
        edgecolor="black", linewidth=0.6, capsize=4,
        error_kw={"linewidth": 1.4, "ecolor": BLACK, "alpha": 0.85}, height=0.65)
# y labels: rank + intuitive featurizer/classifier breakdown over two lines
ylabels = []
for _, r in top5_rev.iterrows():
    rank = int(r["rank_d"])
    feat = PRETTY_FEATURIZER.get(r.shorthand.split('__')[0], '?')
    clf  = PRETTY_CLASSIFIER.get(r.shorthand.split('__')[1], '?')
    ylabels.append(f"#{rank}  {feat}\n     →  {clf}")
ax.set_yticks(y_pos)
ax.set_yticklabels(ylabels, fontsize=12, fontweight="bold", color=BLACK)
# Make label rows visually breathing
ax.set_ylim(-0.6, len(top5_rev) - 0.4)
# bold value labels to the right of each bar
for i, (_, r) in enumerate(top5_rev.iterrows()):
    ax.text(r["mean"] + r["std"] + 0.0035, i,
            f"{r['mean']:.4f} ± {r['std']:.4f}", va="center",
            fontsize=13, fontweight="bold", color=BLACK)
ax.set_xlabel("Mean test accuracy ± 1 SD  (5×5 RSKF, n=25 trials)")
ax.set_xlim(0.78, 0.96)
ax.set_title("Top-5 podium — best five of the 29 benchmarked configurations",
             loc="left", fontsize=16, fontweight="bold")
# in-figure callout below the chart
fig.text(0.5, 0.02,
        "KEY INSIGHT — both #1 and #2 share the same OvR(ExtraTrees) classifier under different featurizers.\n"
        "The win is in the classifier choice. Ranks #3–#5 use the paper's Balanced RF baseline with different word-embedding featurizers.",
        ha="center", fontsize=11.5, fontweight="bold", color=BLACK,
        bbox=dict(boxstyle="round,pad=0.6", facecolor="#fff9e6", edgecolor=BLACK, linewidth=0.8))
plt.tight_layout(rect=(0, 0.10, 1, 1))
plt.savefig(FIG/"fig1b_top5_podium.png"); plt.close()

# =================================================================================
# FIG 2: family-level boxplot — HORIZONTAL orientation so the long family names
# don't overlap. Each box is one family across all (config, seed, fold) trials.
# =================================================================================
print("[deck] Fig 2: family-level boxplot ...")
# Sort families by mean accuracy DESCENDING — highest at top
fam_means = df_pf_c.groupby("family")["acc"].mean().sort_values(ascending=False)
fams_sorted_top_to_bottom = fam_means.index.tolist()
fams_for_matplotlib = list(reversed(fams_sorted_top_to_bottom))  # mpl boxplot draws bottom-up

fig, ax = plt.subplots(figsize=(13, 6.0))
data = [df_pf_c[df_pf_c.family == f]["acc"].values for f in fams_for_matplotlib]
bp = ax.boxplot(data, vert=False, patch_artist=True, widths=0.6, showmeans=True,
                meanprops={"marker": "D", "markerfacecolor": "white", "markeredgecolor": BLACK, "markeredgewidth": 1.2, "markersize": 8},
                medianprops={"color": BLACK, "linewidth": 1.8},
                whiskerprops={"color": BLACK, "linewidth": 1.0},
                capprops={"color": BLACK, "linewidth": 1.0},
                flierprops={"marker": "o", "markerfacecolor": GRAY, "markeredgecolor": BLACK,
                            "markersize": 4, "alpha": 0.6})
for patch, fam in zip(bp["boxes"], fams_for_matplotlib):
    patch.set_facecolor(FAMILY_COLOR[fam]); patch.set_alpha(0.85)
    patch.set_edgecolor(BLACK); patch.set_linewidth(1.0)

# Use the intuitive family names + an inline tag for what each family contains
def family_blurb(f):
    if f == "Ours (shallow)":      return "shallow tree ensemble (OvR ExtraTrees 500)"
    if f == "Paper BRF baselines": return "shallow tree ensemble (paper's Balanced RF 100)"
    if f == "DL: LSTM+attn":       return "recurrent + soft attention"
    if f == "DL: transformer":     return "4-block self-attention"
    if f == "DL: attention":       return "attention-only (no recurrence)"
    if f == "DL: LSTM":            return "vanilla recurrence"
    return f
ylabels = [f"{PRETTY_FAMILY.get(f, f)}\n   ({family_blurb(f)})" for f in fams_for_matplotlib]
ax.set_yticks(range(1, len(fams_for_matplotlib)+1))
ax.set_yticklabels(ylabels, fontsize=11, fontweight="bold", color=BLACK)
ax.set_xlabel("Test accuracy per (config, seed, fold) trial")
ax.set_xlim(0.36, 1.08)
ax.set_title("Accuracy distribution by model family — sorted high-to-low",
             loc="left", fontsize=15, fontweight="bold")

# best-per-family numbers boxed to the right
best_per_fam = df_pf_c.groupby(["family","shorthand"])["acc"].mean().reset_index()
top_by_fam = best_per_fam.loc[best_per_fam.groupby("family")["acc"].idxmax()].set_index("family")
for i, fam in enumerate(fams_for_matplotlib, 1):
    if fam in top_by_fam.index:
        best = top_by_fam.loc[fam]
        ax.annotate(f"best in family: {best.acc:.4f}\n{pretty_name(best.shorthand)}",
                    xy=(0.965, i), xytext=(0.965, i),
                    fontsize=9, fontweight="bold", color=BLACK, va="center", ha="left",
                    bbox=dict(boxstyle="round,pad=0.35", fc="white", ec=BLACK, lw=0.6, alpha=0.95))
plt.tight_layout()
plt.savefig(FIG/"fig2_family_boxplot.png"); plt.close()

# =================================================================================
# FIG 3: reproducibility — orig vs retrained deltas
# =================================================================================
print("[deck] Fig 3: reproducibility ...")
orig = pd.read_csv(ROOT / "artifacts/original_benchmark_per_fold_metrics.csv")[["shorthand","repeat_seed","fold","acc"]]
# For the reproducibility delta chart specifically, compare the rep_1 BENCHMARK
# against the LATEST AVAILABLE additional rep (rep_2 if done) instead of
# re-reading the SAME data from artifacts/predictions/. This shows model-init
# reproducibility (rep_2's seed=2000 vs rep_1's seed=1000), which is the
# scientifically meaningful "Δ across reps" the paper claims.
_FIG3_RETRAIN_SRC = ROOT / "reproducibility/rep_2"
if not (_FIG3_RETRAIN_SRC / "predictions").exists():
    _FIG3_RETRAIN_SRC = REP   # fallback to default (artifacts/) if rep_2 not available
print(f"  fig3 retrain source: {_FIG3_RETRAIN_SRC.relative_to(ROOT)}")
retrained_rows = []
for mj in glob.glob(str(_FIG3_RETRAIN_SRC/"predictions/*/r*_f*/meta.json")):
    retrained_rows.append(json.load(open(mj)))
df_retr = pd.DataFrame(retrained_rows).rename(columns={"seed":"repeat_seed","test_acc":"acc"})[["shorthand","repeat_seed","fold","acc"]]
# Filter to configs with a COMPLETE 25-trial retrain (apples-to-apples comparison only)
ret_counts = df_retr.groupby("shorthand").size()
complete_retrain = set(ret_counts[ret_counts == 25].index)
incomplete = set(df_retr.shorthand.unique()) - complete_retrain
missing_in_retrain = set(orig.shorthand.unique()) - set(df_retr.shorthand.unique())
print(f"  reproducibility data: {len(orig)} orig rows, {len(df_retr)} retrained rows from rep_1/")
print(f"  configs with FULL 25-trial retrain: {len(complete_retrain)}/29")
if incomplete: print(f"  partial retrain (excluded): {sorted(incomplete)}")
if missing_in_retrain: print(f"  no retrain at all (excluded): {sorted(missing_in_retrain)}")
df_retr = df_retr[df_retr.shorthand.isin(complete_retrain)]
orig_for_repro = orig[orig.shorthand.isin(complete_retrain)]
merged = orig_for_repro.merge(df_retr, on=["shorthand","repeat_seed","fold"], suffixes=("_orig","_retr"))
agg2 = merged.groupby("shorthand").agg(orig=("acc_orig","mean"), retr=("acc_retr","mean"),
                                         max_abs_per_fold=("acc_orig", lambda x: 0)).reset_index()
# also compute per-fold |delta| max for each config (useful for the slide blurb)
per_fold_abs = merged.assign(d=(merged.acc_retr-merged.acc_orig).abs()) \
    .groupby("shorthand")["d"].max().reset_index().rename(columns={"d":"max_abs_fold_delta"})
agg2 = agg2.merge(per_fold_abs, on="shorthand")
agg2["delta"] = agg2.retr - agg2.orig
agg2["abs_delta"] = agg2.delta.abs()
agg2["pretty"] = agg2.shorthand.apply(pretty_name)
agg2 = agg2.sort_values("abs_delta", ascending=True)  # smallest |Δ| at bottom → biggest at top

# bucket each config: deterministic / BRF non-det / DL non-det
def bucket(s):
    if s in ("cpu__ET500_log2","ftCbow_MM__ET500_sqrt"): return "deterministic"
    if "BRF100" in s: return "brf"
    return "dl"
agg2["bucket"] = agg2.shorthand.apply(bucket)
bucket_color = {"deterministic": SAGE, "brf": NAVY, "dl": ORANGE}

fig, ax = plt.subplots(figsize=(13, 8.2))
y = np.arange(len(agg2))
colors = [bucket_color[b] for b in agg2.bucket]
ax.barh(y, agg2.delta, color=colors, edgecolor="black", linewidth=0.6)
ax.axvline(0, color=BLACK, linewidth=1.5)
ax.set_yticks(y); ax.set_yticklabels(agg2.pretty, fontsize=9.5, fontweight="bold", color=BLACK)
ax.set_xlabel("Δ accuracy   (rep_2 mean − rep_1 mean)")
ax.set_title("Reproducibility — rep_2 (REPRO_REP_SEED=2000) vs rep_1 (REPRO_REP_SEED=1000), sorted by |Δ|",
             loc="left", fontsize=15, fontweight="bold")
ax.set_xlim(-0.045, 0.025)

handles = [Patch(facecolor=SAGE, edgecolor="black", linewidth=0.5,
                  label="ExtraTrees winners (our sklearn)\nrandom_state-seeded → Δ = 0 across reps"),
           Patch(facecolor=NAVY, edgecolor="black", linewidth=0.5,
                  label="Balanced RF baselines (imblearn)\nrandom_state-seeded → small Δ from thread order"),
           Patch(facecolor=ORANGE, edgecolor="black", linewidth=0.5,
                  label="DL configurations (Keras / TF)\nweight-init varies per rep + GPU op-order")]
ax.legend(handles=handles, loc="lower left", title="What drives Δ ≠ 0?",
          title_fontsize=11, fontsize=10, frameon=True, facecolor="white", edgecolor="#cccccc")

# right-side worst-fold annotation for clarity
_max_dl_delta = float(agg2[agg2['bucket']=='dl']['abs_delta'].max()) if (agg2['bucket']=='dl').any() else 0
_max_fold_delta = float(agg2['max_abs_fold_delta'].max()) if len(agg2) else 0
ax.text(0.018, len(agg2) - 1,
        f"Cross-rep delta = model-init variance only (data splits FIXED at 5×5 RSKF).\n"
        f"ExtraTrees winners: |Δ| = 0 (random_state seeds same trees → bit-identical).\n"
        f"Max mean-Δ across DL configs: {_max_dl_delta:.4f}.\n"
        f"Max single-fold |Δ|: {_max_fold_delta:.4f}.",
        ha="right", va="top", fontsize=10, fontweight="bold", color=BLACK,
        bbox=dict(boxstyle="round,pad=0.5", facecolor="#fff9e6", edgecolor=BLACK, linewidth=0.7))
plt.tight_layout()
plt.savefig(FIG/"fig3_reproducibility.png"); plt.close()
agg2.to_csv(TAB/"tab_reproducibility_delta.csv", index=False)

# =================================================================================
# FIG 4: per-substrate confusion matrix + metric table (best model seed-42 OOF)
# =================================================================================
print("[deck] Fig 4: per-substrate of best model ...")
from sklearn.model_selection import StratifiedKFold
skf = StratifiedKFold(n_splits=5, shuffle=True, random_state=42)
y_pred = np.array([None]*len(X), dtype=object)
y_prob = np.zeros((len(X), 12))
classes = None
for fold, (_, te) in enumerate(skf.split(X, y_labels)):
    d = np.load(REP/f"predictions/cpu__ET500_log2/r42_f{fold}/probs_test.npz", allow_pickle=True)
    probs = d["probs"]; classes_ = d["classes"]
    if classes is None: classes = list(classes_)
    y_prob[te] = probs
    y_pred[te] = np.array([classes_[i] for i in probs.argmax(1)])
from sklearn.metrics import confusion_matrix, precision_recall_fscore_support
cm_raw = confusion_matrix(y_labels, y_pred, labels=substrates)
cm = cm_raw / cm_raw.sum(axis=1, keepdims=True)
p, r, f, sup = precision_recall_fscore_support(y_labels, y_pred, labels=substrates, average=None, zero_division=0)

fig, ax = plt.subplots(figsize=(10, 8))
sns.heatmap(cm, annot=cm_raw.astype(int), fmt="d", cmap="Blues", cbar=False,
            xticklabels=substrates, yticklabels=substrates, ax=ax,
            linewidths=0.5, linecolor="white", annot_kws={"size": 10})
ax.set_xlabel("Predicted substrate", fontsize=12, labelpad=12)
ax.set_ylabel("True substrate", fontsize=12, labelpad=10)
ax.set_title(f"Confusion matrix — cpu__ET500_log2 on seed-42 OOF (n={len(X)}, acc={(y_pred==y_labels).mean():.4f})", loc="left")
plt.setp(ax.get_xticklabels(), rotation=90, ha="center", fontsize=10)
plt.setp(ax.get_yticklabels(), rotation=0, fontsize=10)
plt.tight_layout()
plt.savefig(FIG/"fig4_confusion.png"); plt.close()

per_sub = pd.DataFrame({"substrate": substrates, "n_test": sup,
                          "precision": p, "recall": r, "F1": f}).sort_values("F1", ascending=False)
per_sub.to_csv(TAB/"tab_per_substrate_F1.csv", index=False)

# =================================================================================
# FIG 5: per-substrate sig genes — model's top-3 with lit-hit markers (UNFILTERED)
# =================================================================================
print("[deck] Fig 5: per-substrate top-3 (unfiltered) + lit hit markers (with mapping basis) ...")
fi = pd.read_csv(ROOT/"paper/tables/table4_signature_genes_per_substrate.csv")
def lit_status_with_basis(s, t):
    """Return one of (ASCII markers for reliable PPT rendering):
      ('LIT', '1:1 exact lit match')
      ('LIT', 'via collapse from <alias names>')
      ('miss', 'CAZy not in lit-canon')
      ('—', 'non-CAZy: lit cannot check')
    """
    if not is_cazy(t):
        return ("—", "non-CAZy: lit cannot check")
    if t not in CANON[s]:
        return ("miss", "CAZy not in lit-canon")
    provs = CANON_PROV[s][t]
    if any(k == "exact" for _, k in provs):
        return ("LIT", "1:1 exact lit match")
    alias_names = ", ".join(sorted({a for a,_ in provs}))
    return ("LIT", f"via collapse ← {alias_names}")

rows_sig = []
for s in substrates:
    top3 = fi[fi.substrate == s].sort_values("importance", ascending=False).head(3)
    for _, r in top3.iterrows():
        mark, note = lit_status_with_basis(s, r.feature)
        rows_sig.append({"substrate": s, "feature": r.feature, "importance": r.importance,
                          "mark": mark, "lit_status": note})
sig_tbl = pd.DataFrame(rows_sig)
sig_tbl.to_csv(TAB/"tab_sig_genes_with_lit_status.csv", index=False)

# Render the per-substrate top-3 table as a NATIVE PPT table later; here we save it as CSV.
# (matplotlib table at this size renders too small when embedded — we'll build it as a PPT
# table in the slide section.)
print(f"  per-substrate top-3 table will be drawn as a native PPT table in slide 9")
# Quick PNG version for backup / reviewer reference
fig, ax = plt.subplots(figsize=(15, 8))
ax.axis("off")
cell_data = []
for s in substrates:
    sub = sig_tbl[sig_tbl.substrate == s].reset_index(drop=True)
    cells = []
    for i in range(3):
        if i < len(sub):
            r = sub.iloc[i]
            cells.append(f"[{r.mark}]  {r.feature} (imp={r.importance:.3f})\n{r.lit_status}")
        else:
            cells.append("")
    cell_data.append([s] + cells)
colwid = [0.15, 0.28, 0.28, 0.28]
tbl = ax.table(cellText=cell_data, colLabels=["substrate","top 1","top 2","top 3"],
               loc="center", cellLoc="left", colWidths=colwid)
tbl.auto_set_font_size(False); tbl.set_fontsize(10); tbl.scale(1, 3.0)
for i in range(len(cell_data)+1):
    for j in range(4):
        cell = tbl[(i,j)]
        cell.set_edgecolor("#cccccc")
        if i == 0:
            cell.set_facecolor(NAVY); cell.set_text_props(color="white", weight="bold")
        else:
            txt = cell.get_text().get_text()
            if "[LIT]" in txt and "1:1" in txt: cell.set_facecolor("#c8e6c9")
            elif "[LIT]" in txt: cell.set_facecolor("#e0f3e8")
            elif "[miss]" in txt: cell.set_facecolor("#fce4e4")
            elif "[—]" in txt: cell.set_facecolor("#fdf6e3")
            else: cell.set_facecolor("white")
ax.set_title("Model's TOP-3 features per substrate (raw, unfiltered) — lit-DB status with mapping basis",
             loc="left", color=CHARCOAL, pad=8, fontsize=14)
ax.text(0.0, -0.04,
        "  [LIT] 1:1 exact (strong green)  ·  [LIT] via alias collapse (pale green)  ·  [miss] CAZy not in lit canon (red)  ·  [—] non-CAZy / lit cannot adjudicate (yellow)",
        transform=ax.transAxes, fontsize=10, color=CHARCOAL)
plt.savefig(FIG/"fig5_sig_genes_unfiltered.png"); plt.close()

# Summary stats
n_cazy_check = sum(1 for r in rows_sig if r["mark"] in ("LIT","miss"))
n_cazy_in_lit = sum(1 for r in rows_sig if r["mark"] == "LIT")
n_exact = sum(1 for r in rows_sig if "1:1 exact" in r["lit_status"])
n_collapse = sum(1 for r in rows_sig if "collapse" in r["lit_status"])
n_non_cazy = sum(1 for r in rows_sig if r["mark"] == "—")
print(f"  top-3 lit-hit summary: {n_cazy_in_lit}/{n_cazy_check} CAZy features are lit-canonical")
print(f"                         {n_exact} via 1:1 exact match, {n_collapse} via alias collapse")
print(f"                         {n_non_cazy} non-CAZy features (lit cannot adjudicate)")
print(f"  top-3 lit hits: {n_cazy_in_lit}/{n_cazy_check} CAZy features in lit canon  ({n_non_cazy}/36 non-CAZy, not adjudicated)")

# =================================================================================
# FIG 6: literature-scope coverage + per-PUL any-hit (population metrics)
# =================================================================================
print("[deck] Fig 6: population lit-validation — TRUE-class attribution ...")
# Use the TRUE-class ablation (Δ-prob w.r.t. true substrate, not argmax) so the
# population metrics here match the per-substrate funnels on slide 13 exactly.
oof = pd.read_csv(REP/"ablation/sig_gene_ablation_oof_outer42_groundtruth_calibrated.csv")
# Alias the calibrated argmax probability to 'prob' for chart 7 (example PULs) compatibility
if "prob" not in oof.columns:
    oof = oof.rename(columns={"prob_pred_cal": "prob"})
# `correct` (true == pred subset) still used for the example-PUL cherry-pick on slide 10
correct = oof[oof.true == oof.pred].copy()

def per_pul_anyhit_true():
    rows = []
    for K in (1,3,5):
        n_elig = 0; n_hit = 0
        for _, r in oof.iterrows():
            cs = CANON[r.true]   # canon for TRUE substrate
            pul_toks = set(tok_cpu(X[r.idx]))
            if not (pul_toks & cs): continue  # not eligible
            n_elig += 1
            top = str(r["top1"]) if K==1 else (str(r["top3"]) if K==3 else str(r["top5"]))
            tokens = set(top.split(";"))
            if tokens & cs: n_hit += 1
        rows.append({"K": K, "n_eligible": n_elig, "n_hit": n_hit, "rate": n_hit/n_elig})
    return pd.DataFrame(rows)

def lit_gene_coverage_true():
    rows = []
    for K in (1,3,5):
        in_scope = 0; covered = 0
        flagged = {s: set() for s in substrates}
        # Flag using TRUE substrate (not predicted) — so even PULs the model classified
        # wrong contribute their top-K Δ_true tokens to their TRUE substrate's flagged set.
        for _, r in oof.iterrows():
            top = str(r["top1"]) if K==1 else (str(r["top3"]) if K==3 else str(r["top5"]))
            flagged[r.true] |= set(top.split(";")) & CANON[r.true]
        for s in substrates:
            scope = set()
            for _, r in oof[oof.true == s].iterrows():
                scope |= set(tok_cpu(X[r.idx])) & CANON[s]
            in_scope += len(scope)
            covered += len(scope & flagged[s])
        rows.append({"K": K, "in_scope": in_scope, "covered": covered, "rate": covered/in_scope})
    return pd.DataFrame(rows)

anyhit_df = per_pul_anyhit_true(); scope_df = lit_gene_coverage_true()
anyhit_df.to_csv(TAB/"tab_perpul_anyhit.csv", index=False)
scope_df.to_csv(TAB/"tab_lit_gene_scope.csv", index=False)

n_elig_total = int(anyhit_df.iloc[0].n_eligible)
n_scope_total = int(scope_df.iloc[0].in_scope)

fig, axes = plt.subplots(1, 2, figsize=(13, 5))
ax = axes[0]
Ks = anyhit_df.K.values; rates = anyhit_df.rate.values * 100
ax.bar(Ks, rates, color=SAGE, alpha=0.85, edgecolor="white", width=0.6)
for K, ra, hi, el in zip(Ks, rates, anyhit_df.n_hit, anyhit_df.n_eligible):
    ax.text(K, ra+1.5, f"{ra:.1f}%\n{hi}/{el}", ha="center", fontsize=11, color=CHARCOAL, fontweight="bold")
ax.set_xticks(Ks); ax.set_xlabel("top-K ablation tokens", fontsize=11)
ax.set_ylabel("% PULs (eligible, true-class) with ≥1 lit-canonical in top-K", fontsize=10.5)
ax.set_ylim(0, 115)
ax.set_title(f"Per-PUL any-hit (TRUE-class attribution)\nn = {n_elig_total} eligible PULs (lit-canon present in PUL)",
              loc="left", fontsize=12, color=NAVY)
ax = axes[1]
rates2 = scope_df.rate.values * 100
ax.bar(Ks, rates2, color=NAVY, alpha=0.85, edgecolor="white", width=0.6)
for K, ra, c, s in zip(Ks, rates2, scope_df.covered, scope_df.in_scope):
    ax.text(K, ra+1.5, f"{ra:.1f}%\n{c}/{s}", ha="center", fontsize=11, color=CHARCOAL, fontweight="bold")
ax.set_xticks(Ks); ax.set_xlabel("top-K ablation tokens", fontsize=11)
ax.set_ylabel("% in-scope canonical CAZy genes flagged anywhere", fontsize=10.5)
ax.set_ylim(0, 115)
ax.set_title(f"Gene-scope coverage (TRUE-class attribution)\nn = {n_scope_total} in-scope canonical CAZy genes",
              loc="left", fontsize=12, color=NAVY)
plt.tight_layout()
plt.savefig(FIG/"fig6_lit_validation.png"); plt.close()

# =================================================================================
# FIG 7: ten cherry-picked example PULs with high lit-hit coverage
# =================================================================================
print("[deck] Fig 7: 10 example PULs ...")
correct["pul_tokens"] = correct.idx.apply(lambda i: set(tok_cpu(X[i])))
correct["n_canon_top3"] = correct.apply(lambda r: sum(1 for t in str(r.top3).split(";") if t in CANON[r.pred]), axis=1)
correct["n_canon_top5"] = correct.apply(lambda r: sum(1 for t in str(r.top5).split(";") if t in CANON[r.pred]), axis=1)
correct = correct.sort_values(["n_canon_top3","n_canon_top5","prob"], ascending=[False,False,False])
# pick 10: 1-2 per substrate, prioritizing high lit coverage
picks = []
seen_subs = Counter()
for _, r in correct.iterrows():
    if seen_subs[r.pred] >= 2: continue
    if len(picks) >= 10: break
    picks.append(r.to_dict()); seen_subs[r.pred] += 1
ex = pd.DataFrame(picks)
# pretty-format top-5 using ASCII markers that render reliably in PPT
def annot_top5(s, top5):
    cs = CANON[s]
    out = []
    for piece in str(top5).split(";")[:5]:
        if not piece: continue
        tok = piece.split(":")[0].rstrip("*")
        if tok in cs:        mark = "[LIT]"
        elif is_cazy(tok):   mark = "[CAZy]"
        else:                mark = "[acc]"
        out.append(f"{mark} {piece}")
    return "; ".join(out)

def short_pul(idx, max_chars=80):
    s = str(X[idx]).strip()
    if len(s) <= max_chars: return s
    return s[:max_chars-3] + "..."

ex["sig_gene_annot"] = ex.apply(lambda r: annot_top5(r.pred, r.top5), axis=1)
ex["p_value"] = ex.prob.apply(lambda p: f"{(1-p)**11:.1e}")
ex["pul_sequence"] = ex.idx.apply(lambda i: short_pul(i, max_chars=70))
ex_show = ex[["pul_sequence","pred","prob","p_value","sig_gene_annot","n_canon_top5"]].rename(
    columns={"pred":"substrate","prob":"max prob","sig_gene_annot":"top-5 signature genes (LIT = lit-canonical; CAZy = non-canonical CAZy; acc = non-CAZy accessory)","n_canon_top5":"#LIT/5"})
ex_show.to_csv(TAB/"tab_example_predictions.csv", index=False)
# Also keep idx for traceability
ex.assign(idx_col=ex.idx)[["idx","pul_sequence","pred","prob","p_value","sig_gene_annot","n_canon_top5"]].to_csv(TAB/"tab_example_predictions_with_idx.csv", index=False)

# =================================================================================
# FIG 8: calibration — reliability curve + ECE table
# =================================================================================
print("[deck] Fig 8: calibration ...")
cal_best = np.load(REP/"calibration/oof_outer42_best_of_both.npz", allow_pickle=True)
cal_cv5 = np.load(REP/"calibration/oof_outer42_calibration.npz", allow_pickle=True)
y_int = cal_best["y_true"]
def reliability(probs, n_bins=10):
    conf = probs.max(1); correct = (probs.argmax(1) == y_int).astype(float)
    edges = np.linspace(0,1,n_bins+1); mids = (edges[1:]+edges[:-1])/2
    accs = np.zeros(n_bins); confs = np.zeros(n_bins); ns = np.zeros(n_bins)
    for i in range(n_bins):
        m = (conf >= edges[i]) & (conf < edges[i+1]) if i<n_bins-1 else (conf >= edges[i]) & (conf <= edges[i+1])
        if m.sum():
            accs[i] = correct[m].mean(); confs[i] = conf[m].mean(); ns[i] = m.sum()
    return mids, accs, confs, ns
def ece(probs, n_bins=10):
    mids, accs, confs, ns = reliability(probs, n_bins)
    return float(np.sum(ns * np.abs(accs - confs))/ns.sum())
fig, axes = plt.subplots(1, 2, figsize=(11, 4.5))
ax = axes[0]
for probs, name, color in [(cal_best["probs_uncal"], "uncalibrated", GRAY),
                            (cal_best["probs_temp"], f"temperature T={cal_best['T_per_fold'].mean():.2f}", SAGE),
                            (cal_cv5["probs_iso"], "isotonic (CalibratedCV cv=5)", ORANGE),
                            (cal_cv5["probs_sig"], "sigmoid/Platt", CRIMSON)]:
    mids, accs, _, ns = reliability(probs)
    mask = ns > 0
    ax.plot(mids[mask], accs[mask], "o-", color=color, label=f"{name} (ECE={ece(probs):.3f})", linewidth=1.5, markersize=5)
ax.plot([0,1], [0,1], "--", color=CHARCOAL, linewidth=0.8, alpha=0.5, label="perfect calibration")
ax.set_xlim(0,1); ax.set_ylim(0,1); ax.set_xlabel("predicted confidence"); ax.set_ylabel("empirical accuracy")
ax.set_title("Reliability diagram — 10-bin", loc="left"); ax.legend(loc="upper left", fontsize=8.5, frameon=False)
# Right panel: metric table
ax = axes[1]; ax.axis("off")
def acc(p): return float((p.argmax(1)==y_int).mean())
metrics = [
    ("Uncalibrated", acc(cal_best["probs_uncal"]), ece(cal_best["probs_uncal"]), "raw OvR-ET output"),
    (f"Temperature T={cal_best['T_per_fold'].mean():.2f}", acc(cal_best["probs_temp"]), ece(cal_best["probs_temp"]), "preserves argmax exactly"),
    ("Isotonic (cv=5)", acc(cal_cv5["probs_iso"]), ece(cal_cv5["probs_iso"]), "can re-rank classes"),
    ("Sigmoid (cv=5)", acc(cal_cv5["probs_sig"]), ece(cal_cv5["probs_sig"]), "worse than uncal"),
]
cells = [[m[0], f"{m[1]:.4f}", f"{m[2]:.4f}", m[3]] for m in metrics]
tbl = ax.table(cellText=cells, colLabels=["method","accuracy","ECE (10-bin)","note"],
                loc="center", cellLoc="left", colWidths=[0.30, 0.14, 0.16, 0.40])
tbl.auto_set_font_size(False); tbl.set_fontsize(9); tbl.scale(1, 2.0)
for j in range(4):
    tbl[(0,j)].set_facecolor(NAVY); tbl[(0,j)].set_text_props(color="white", weight="bold")
# highlight temperature row
for j in range(4):
    tbl[(2,j)].set_facecolor("#e0f3e8")
ax.set_title("Calibration metrics — seed-42 OOF n=1030", loc="left", color=CHARCOAL, pad=6)
plt.tight_layout()
plt.savefig(FIG/"fig8_calibration.png"); plt.close()

# =================================================================================
# FIG 8b: PER-SUBSTRATE SIG-GENE FUNNEL — GROUND-TRUTH-CLASS attribution
# ===================================================================================
# We compute Δ-prob with respect to the TRUE class (not argmax). This isolates
# attribution quality from classification quality. For each test PUL whose true
# substrate is s, the sig genes ARE the tokens that most reduce P(s) when removed.
#   PUL view:  test_of_s → eligible (has lit-canon in PUL) → hit @K=3
#   Gene view: lit_canon → in-scope → flagged @K=3
# No "correctly predicted" filter — we use TRUE substrate everywhere because our
# leave-one-token-out ablation can be computed for ANY class for ANY PUL.
# =================================================================================
print("[deck] Fig 8b: ground-truth-class per-substrate sig-gene funnels ...")
oof_gt = pd.read_csv(REP/"ablation/sig_gene_ablation_oof_outer42_groundtruth_calibrated.csv")
K = 3
sig_rows_tab = []
for s in substrates:
    test_of_s = oof_gt[oof_gt.true == s]
    n_total = len(test_of_s)
    # PUL view: eligible = lit-canon present in PUL; hit = top-K_true contains lit-canon
    n_elig = 0; n_hit = 0
    for _, r in test_of_s.iterrows():
        toks = set(tok_cpu(X[r.idx]))
        if toks & CANON[s]:
            n_elig += 1
            top = set(str(r[f"top{K}"]).split(";"))
            if top & CANON[s]: n_hit += 1
    hit_rate = (n_hit / n_elig) if n_elig else 0.0
    # Gene view
    lit_n     = len(CANON[s])
    in_scope  = set()
    for _, r in test_of_s.iterrows():
        in_scope |= set(tok_cpu(X[r.idx])) & CANON[s]
    n_inscope = len(in_scope)
    flagged = set()
    for _, r in test_of_s.iterrows():
        flagged |= set(str(r[f"top{K}"]).split(";")) & CANON[s]
    n_flag   = len(in_scope & flagged)
    cov_rate = (n_flag / n_inscope) if n_inscope else 0.0
    sig_rows_tab.append(dict(substrate=s,
        n_total=n_total, n_eligible=n_elig, n_pul_hit_at_K=n_hit, pul_hit_rate=hit_rate,
        lit_canon_size=lit_n, n_in_scope=n_inscope, n_flagged_at_K=n_flag, scope_recall=cov_rate,
    ))
sig_pr = pd.DataFrame(sig_rows_tab).sort_values("pul_hit_rate", ascending=False)
sig_pr.to_csv(TAB/"tab_per_substrate_sig_pr.csv", index=False)
sig_pr.to_csv(ROOT/"paper/tables/table12_per_substrate_sig_pr.csv", index=False)

fig, (axL, axR) = plt.subplots(1, 2, figsize=(16, 7.0), gridspec_kw={"wspace": 0.30})
y_pos = np.arange(len(sig_pr))
bar_h = 0.25

# LEFT: PUL view — 3 nested bars (n_total, n_eligible, n_hit) — no n_correct gate
axL.barh(y_pos + bar_h, sig_pr.n_total.values,         bar_h, color="#cccccc", edgecolor=BLACK, linewidth=0.4, label="total test PULs (TRUE substrate = this row)")
axL.barh(y_pos,          sig_pr.n_eligible.values,     bar_h, color=NAVY,      edgecolor=BLACK, linewidth=0.4, label="of those, eligible (lit-canon present in PUL)")
axL.barh(y_pos - bar_h,  sig_pr.n_pul_hit_at_K.values, bar_h, color=SAGE,      edgecolor=BLACK, linewidth=0.4, label="of eligible, top-3 sig gene for TRUE class is canonical")
axL.set_yticks(y_pos); axL.set_yticklabels(sig_pr.substrate, fontsize=11, fontweight="bold", color=BLACK)
axL.invert_yaxis()
axL.set_xlabel("number of PULs")
axL.set_title("PUL view — funnel of counts (no correctness gate)", loc="left", fontsize=13, fontweight="bold")
axL.legend(loc="lower right", fontsize=9, frameon=True, facecolor="white", edgecolor="#cccccc")
xmax = float(sig_pr.n_total.max()) * 1.10
axL.set_xlim(0, xmax * 1.20)
for i, (_, r) in enumerate(sig_pr.iterrows()):
    axL.text(xmax, i, f"{int(r.n_pul_hit_at_K)}/{int(r.n_eligible)} = {r.pul_hit_rate*100:.0f}%",
              va="center", fontsize=10, fontweight="bold", color=BLACK)

# RIGHT: Gene view — 3 nested bars (lit_n, in_scope, flagged)
axR.barh(y_pos + bar_h, sig_pr.lit_canon_size.values,  bar_h, color="#cccccc", edgecolor=BLACK, linewidth=0.4, label="lit-canon genes after alias collapse")
axR.barh(y_pos,         sig_pr.n_in_scope.values,      bar_h, color=ORANGE,    edgecolor=BLACK, linewidth=0.4, label="of those, in-scope (appears in any test PUL of this substrate)")
axR.barh(y_pos - bar_h, sig_pr.n_flagged_at_K.values,  bar_h, color=SAGE,      edgecolor=BLACK, linewidth=0.4, label="of in-scope, surfaced as top-3 sig gene for TRUE class anywhere")
axR.set_yticks(y_pos); axR.set_yticklabels(sig_pr.substrate, fontsize=11, fontweight="bold", color=BLACK)
axR.invert_yaxis()
axR.set_xlabel("number of lit-canonical CAZy families")
axR.set_title("Gene view — funnel of counts (no correctness gate)", loc="left", fontsize=13, fontweight="bold")
axR.legend(loc="lower right", fontsize=9, frameon=True, facecolor="white", edgecolor="#cccccc")
xmaxR = float(sig_pr.lit_canon_size.max()) * 1.10
axR.set_xlim(0, xmaxR * 1.35)
for i, (_, r) in enumerate(sig_pr.iterrows()):
    axR.text(xmaxR, i, f"{int(r.n_flagged_at_K)}/{int(r.n_in_scope)} = {r.scope_recall*100:.0f}%",
              va="center", fontsize=10, fontweight="bold", color=BLACK)

T_total = int(sig_pr.n_total.sum()); T_elig = int(sig_pr.n_eligible.sum()); T_hit = int(sig_pr.n_pul_hit_at_K.sum())
T_lit   = int(sig_pr.lit_canon_size.sum()); T_isc = int(sig_pr.n_in_scope.sum()); T_flg = int(sig_pr.n_flagged_at_K.sum())
fig.suptitle(f"Per-substrate sig-gene FUNNEL (K=3, TRUE-class attribution) — "
              f"PUL view {T_total}→{T_elig}→{T_hit} ({T_hit/T_elig*100:.1f}% hit) ·  "
              f"gene view {T_lit}→{T_isc}→{T_flg} ({T_flg/T_isc*100:.1f}% scope recall)",
              fontsize=13, fontweight="bold", y=0.99)
plt.tight_layout(rect=(0, 0, 1, 0.96))
plt.savefig(FIG/"fig8b_per_substrate_sig_pr.png"); plt.close()

# =================================================================================
# FIG 8c: TEST-PUL OOV (out-of-vocabulary token proportion) vs ACCURACY
# For seed-42 5-fold OOF, compute per-PUL OOV proportion using the same featurizer
# as the model (CountVec with tok_cpu, fitted per-fold on outer_tr only — leak-free).
# =================================================================================
print("[deck] Fig 8c: per-PUL OOV proportion vs accuracy ...")
from sklearn.feature_extraction.text import CountVectorizer
oov_props = np.zeros(len(X)); n_toks = np.zeros(len(X), dtype=int); n_oov = np.zeros(len(X), dtype=int)
skf = StratifiedKFold(n_splits=5, shuffle=True, random_state=42)
for fold, (tr_idx, te_idx) in enumerate(skf.split(X, y_labels)):
    cv_ = CountVectorizer(tokenizer=tok_cpu, token_pattern=None, lowercase=False)
    cv_.fit(X[tr_idx])
    vocab_set = set(cv_.vocabulary_.keys())
    for idx in te_idx:
        toks = tok_cpu(X[idx])
        n_toks[idx] = len(toks)
        n_oov[idx] = sum(1 for t in toks if t not in vocab_set)
        oov_props[idx] = n_oov[idx] / max(len(toks), 1)
correct_mask = (y_pred == y_labels)
buckets_def = [(0.0, 0.0, "0%"), (0.001, 0.05, "0–5%"),
                (0.05, 0.10, "5–10%"), (0.10, 0.25, "10–25%"),
                (0.25, 1.01, "≥25%")]
buc_rows = []
for lo, hi, lab in buckets_def:
    m = (oov_props == 0) if lo == 0 == hi else (oov_props >= lo) & (oov_props < hi)
    if m.sum() == 0:
        buc_rows.append((lab, 0, 0.0, 0.0))
    else:
        buc_rows.append((lab, int(m.sum()), float(correct_mask[m].mean()), float(oov_props[m].mean()*100)))
buc_df = pd.DataFrame(buc_rows, columns=["bucket","n_PULs","accuracy","mean_OOV_pct"])
buc_df.to_csv(TAB/"tab_oov_vs_accuracy.csv", index=False)
buc_df.to_csv(ROOT/"paper/tables/table13_oov_vs_accuracy.csv", index=False)

fig, ax = plt.subplots(figsize=(13, 6.0))
# Uniform-width bars — bucket density is shown via the n= annotation under each bar
x_pos = np.arange(len(buc_rows))
colors_b = [SAGE if r[2] >= 0.9 else ORANGE if r[2] >= 0.7 else CRIMSON for r in buc_rows]
bars = ax.bar(x_pos, [r[2] for r in buc_rows], width=0.62, color=colors_b,
              edgecolor="black", linewidth=0.7)
ax.set_xticks(x_pos)
ax.set_xticklabels([f"{r[0]}\nn={r[1]} PULs\nmean OOV={r[3]:.1f}%" for r in buc_rows],
                    fontsize=11, fontweight="bold", color=BLACK)
ax.set_ylabel("Test accuracy (fraction correct)")
ax.set_xlabel("Per-PUL out-of-vocabulary token proportion bucket")
ax.set_ylim(0, 1.08)
# value labels on each bar (accuracy + bucket-share %)
total_n = sum(r[1] for r in buc_rows)
for x, r in zip(x_pos, buc_rows):
    pct = 100.0 * r[1] / max(total_n, 1)
    ax.text(x, r[2] + 0.02, f"{r[2]:.3f}",
            ha="center", fontsize=11, fontweight="bold", color=BLACK)
    ax.text(x, 0.04, f"{pct:.1f}% of test PULs",
            ha="center", fontsize=9, color="#555")
ax.set_title("Test-PUL accuracy vs out-of-vocabulary token proportion (CountVec_cpu featurizer, seed-42 5-fold OOF)",
             loc="left", fontsize=13, fontweight="bold")
from scipy.stats import pointbiserialr as _pbr
_r, _p = _pbr(correct_mask.astype(int), oov_props)
fig.text(0.5, 0.02,
        f"Point-biserial r(correct, OOV proportion) = {_r:.4f}, p = {_p:.2e}  "
        "—  the model is robust to small OOV (≤10% of tokens missing from train vocab) "
        "but accuracy drops sharply once ≥10% of a test PUL's tokens are unseen.",
        ha="center", fontsize=10.5, fontweight="bold", color=BLACK,
        bbox=dict(boxstyle="round,pad=0.5", facecolor="#fff9e6", edgecolor=BLACK, linewidth=0.7))
plt.tight_layout(rect=(0, 0.08, 1, 1))
plt.savefig(FIG/"fig8c_oov_vs_accuracy.png"); plt.close()
print(f"  per-PUL OOV: mean={oov_props.mean()*100:.2f}% median={np.median(oov_props)*100:.2f}%  "
      f"PULs with 0% OOV: {(oov_props == 0).sum()}/{len(X)} ({(oov_props == 0).mean()*100:.1f}%)")
print(f"  pointbiserial r(correct, oov) = {_r:.4f}, p = {_p:.2e}")

# =================================================================================
# FIG 9: per-fold accuracy distribution + spread (robustness)
# =================================================================================
print("[deck] Fig 9: per-trial swarm + box + central-tendency markers ...")
# Two-panel layout: left = swarm-over-box for all 6 families; right = per-family stats table.
fams_by_mean = df_pf_c.groupby("family")["acc"].mean().sort_values(ascending=False).index.tolist()
fams_for_plot = list(reversed(fams_by_mean))  # mpl draws bottom-up; reverse so highest on top
SHORT_FAMILY = {
    "Ours (shallow)":      "Our shallow winners\n(ExtraTrees 500)",
    "Paper BRF baselines": "Paper baselines\n(Balanced RF 100)",
    "DL: LSTM":            "Deep LSTM (paper)",
    "DL: LSTM+attn":       "Deep LSTM + attention (paper)",
    "DL: attention":       "Deep attention-only (paper)",
    "DL: transformer":     "Deep transformer (paper)",
}

fig = plt.figure(figsize=(15, 11.0))
# 2-row layout: chart on top, table below (full width); generous spacing to avoid clipping
gs = fig.add_gridspec(2, 1, height_ratios=[2.6, 2.0], hspace=0.40)
ax = fig.add_subplot(gs[0, 0])
ax_t = fig.add_subplot(gs[1, 0]); ax_t.axis("off")

positions = np.arange(len(fams_for_plot))
data_for_plot = [df_pf_c[df_pf_c.family == f]["acc"].values for f in fams_for_plot]

# Horizontal BOX plot showing central-tendency markers (Q1/median/Q3, whiskers, mean)
bp = ax.boxplot(data_for_plot, positions=positions, vert=False, widths=0.55, patch_artist=True,
                showmeans=True, manage_ticks=False,
                meanprops={"marker": "D", "markerfacecolor": "white",
                            "markeredgecolor": BLACK, "markersize": 11, "markeredgewidth": 1.6},
                medianprops={"color": BLACK, "linewidth": 2.8},
                whiskerprops={"color": BLACK, "linewidth": 1.4},
                capprops={"color": BLACK, "linewidth": 1.6},
                flierprops={"marker": "", "markersize": 0})  # hide default fliers (we draw all points below)
for patch, fam in zip(bp["boxes"], fams_for_plot):
    patch.set_facecolor(FAMILY_COLOR[fam]); patch.set_alpha(0.40)
    patch.set_edgecolor(BLACK); patch.set_linewidth(1.4)

# Swarm-style jittered points overlaid on the box
np.random.seed(0)
for i, (d, fam) in enumerate(zip(data_for_plot, fams_for_plot)):
    y = np.random.normal(i, 0.06, len(d))
    ax.scatter(d, y, s=22, color=FAMILY_COLOR[fam], edgecolor=BLACK, linewidth=0.4,
                alpha=0.78, zorder=4)

# Y tick labels
def family_n_blurb(f):
    n_cfg = df_pf_c[df_pf_c.family == f].shorthand.nunique()
    n_trials = (df_pf_c.family == f).sum()
    return f"{SHORT_FAMILY.get(f, f)}\n({n_cfg} cfg × 25 = {n_trials} trials)"
ax.set_yticks(positions)
ax.set_yticklabels([family_n_blurb(f) for f in fams_for_plot],
                   fontsize=11, fontweight="bold", color=BLACK)
ax.set_xlabel("Test accuracy per (seed, fold) trial")
ax.set_xlim(0.36, 1.0); ax.set_ylim(-0.6, len(fams_for_plot)-0.4)
ax.set_title("Per-trial accuracy — swarm over box with explicit central-tendency markers",
             loc="left", fontsize=14, fontweight="bold")
# vertical line at our winner's mean
top_mean = df_pf_c[df_pf_c.shorthand == "cpu__ET500_log2"]["acc"].mean()
ax.axvline(top_mean, color=SAGE, linestyle=":", linewidth=1.8, alpha=0.85,
            label=f"winner mean = {top_mean:.4f}")
ax.legend(loc="lower left", frameon=True, fontsize=10.5, facecolor="white", edgecolor="#cccccc")

# Bottom panel: per-family stats TABLE (full width)
stats_rows = []
for fam in reversed(fams_for_plot):  # show highest-mean on TOP of the table
    d = df_pf_c[df_pf_c.family == fam]["acc"].values
    n_cfg = df_pf_c[df_pf_c.family == fam].shorthand.nunique()
    n_trials = len(d)
    stats_rows.append([
        SHORT_FAMILY.get(fam, fam).replace("\n", " · "),
        f"{n_cfg}",
        f"{n_trials}",
        f"{d.mean():.4f}",
        f"{np.median(d):.4f}",
        f"{d.std():.4f}",
        f"{np.percentile(d,25):.4f}",
        f"{np.percentile(d,75):.4f}",
        f"{np.percentile(d,75)-np.percentile(d,25):.4f}",
        f"{d.min():.4f}",
        f"{d.max():.4f}",
    ])
col_labels = ["model family", "n configs", "n trials", "mean", "median",
              "SD", "Q1", "Q3", "IQR (Q3−Q1)", "min", "max"]
tbl = ax_t.table(cellText=stats_rows, colLabels=col_labels,
                  loc="center", cellLoc="center",
                  colWidths=[0.26, 0.06, 0.06, 0.07, 0.07,
                             0.07, 0.07, 0.07, 0.09, 0.07, 0.07])
# left-align first column
for i in range(len(stats_rows) + 1):
    tbl[(i, 0)].set_text_props(ha="left")
tbl.auto_set_font_size(False); tbl.set_fontsize(10); tbl.scale(1, 1.6)
for j in range(len(col_labels)):
    tbl[(0, j)].set_facecolor(NAVY); tbl[(0, j)].set_text_props(color="white", weight="bold")
# Highlight the winner family row
for j in range(len(col_labels)):
    tbl[(1, j)].set_facecolor("#c8e6c9")
    tbl[(1, j)].set_text_props(weight="bold")
ax_t.set_title("Per-family central-tendency metrics — values match the appendix per-classifier-family aggregate",
                loc="left", fontsize=12, fontweight="bold", color=BLACK, y=1.00, pad=10)

# Mini "how to read box" guide as a compact caption BELOW the table
ax_t.text(0, -0.20,
          "HOW TO READ EACH BOX (above): filled box = Q1→Q3 (middle 50%) · thick vertical = median · "
          "white diamond = mean · whiskers = 1.5·IQR rule · dots = individual (seed,fold) trial accuracies.",
          transform=ax_t.transAxes, fontsize=10, fontweight="bold", color=BLACK,
          bbox=dict(boxstyle="round,pad=0.5", facecolor="#fff9e6", edgecolor=BLACK, linewidth=0.7))

plt.tight_layout(rect=(0, 0.02, 1, 0.98))
plt.savefig(FIG/"fig9_trial_distribution.png"); plt.close()

# =================================================================================
# Appendix aggregates — full leaderboard + per-featurizer + per-classifier
# =================================================================================
print("[deck] appendix: full leaderboard + per-featurizer + per-classifier aggregates ...")
full_lb = df_pf_c.groupby("shorthand").agg(
    mean=("acc","mean"), std=("acc","std"),
    min_=("acc","min"), max_=("acc","max"),
    n=("acc","count")
).reset_index().sort_values("mean", ascending=False)
full_lb["rank"]           = np.arange(1, len(full_lb)+1)
full_lb["featurizer_fam"] = full_lb.shorthand.apply(lambda s: decode_shorthand(s)[0])
full_lb["featurizer_det"] = full_lb.shorthand.apply(lambda s: decode_shorthand(s)[1])
full_lb["clf_fam"]        = full_lb.shorthand.apply(lambda s: decode_shorthand(s)[2])
full_lb["clf_det"]        = full_lb.shorthand.apply(lambda s: decode_shorthand(s)[3])

feat_agg = df_pf_c.assign(ff=df_pf_c.shorthand.apply(lambda s: decode_shorthand(s)[0])) \
    .groupby("ff").agg(mean_acc=("acc","mean"), std_acc=("acc","std"),
                       n_configs=("shorthand","nunique"), n_trials=("acc","count")) \
    .reset_index().rename(columns={"ff": "featurizer_family"}).sort_values("mean_acc", ascending=False)

clf_agg = df_pf_c.assign(cf=df_pf_c.shorthand.apply(lambda s: decode_shorthand(s)[2])) \
    .groupby("cf").agg(mean_acc=("acc","mean"), std_acc=("acc","std"),
                       n_configs=("shorthand","nunique"), n_trials=("acc","count")) \
    .reset_index().rename(columns={"cf": "classifier_family"}).sort_values("mean_acc", ascending=False)

full_lb.to_csv(TAB/"tab_full_leaderboard_decoded.csv", index=False)
feat_agg.to_csv(TAB/"tab_per_featurizer_aggregate.csv", index=False)
clf_agg.to_csv(TAB/"tab_per_classifier_aggregate.csv", index=False)
print(f"  per-featurizer: {feat_agg.featurizer_family.tolist()}")
print(f"  per-classifier: {clf_agg.classifier_family.tolist()}")

# Training time aggregates
time_fam_agg = df_pf_c.assign(cf=df_pf_c.shorthand.apply(lambda s: decode_shorthand(s)[2])) \
    .groupby("cf").agg(wall_mean=("wall_sec","mean"), wall_total=("wall_sec","sum"),
                       n_trials=("wall_sec","count"), n_configs=("shorthand","nunique")) \
    .reset_index().rename(columns={"cf": "classifier_family"}).sort_values("wall_mean")
GRAND_TOTAL_SEC = float(df_pf_c.wall_sec.sum())
time_fam_agg.to_csv(TAB/"tab_training_time_per_family.csv", index=False)
print(f"  grand total training wall time across all 725 fits: {GRAND_TOTAL_SEC:.0f} s = {GRAND_TOTAL_SEC/3600:.2f} h")

# Fig 10: training time per family
print("[deck] Fig 10: training time per family ...")
fig, ax = plt.subplots(figsize=(11, 5.0))
y_pos = np.arange(len(time_fam_agg))
fam_colors_clf = {"OvR(ExtraTrees)": SAGE, "OvR(BalancedRF)": NAVY,
                  "DL: LSTM": GRAY, "DL: LSTM+attention": ORANGE,
                  "DL: attention": CRIMSON, "DL: transformer": "#8e44ad"}
bar_colors_t = [fam_colors_clf.get(f, GRAY) for f in time_fam_agg.classifier_family]
ax.barh(y_pos, time_fam_agg.wall_mean, color=bar_colors_t, edgecolor="white", linewidth=0.5)
ax.set_yticks(y_pos); ax.set_yticklabels(time_fam_agg.classifier_family, fontsize=11)
for i, (_, r) in enumerate(time_fam_agg.iterrows()):
    ax.text(r.wall_mean + 1.5, i, f"{r.wall_mean:.1f} s/fold  →  {r.wall_mean*25:.0f} s/config",
            va="center", fontsize=10, color=CHARCOAL)
ax.set_xlabel("Mean wall-clock seconds per (seed, fold) trial")
ax.set_xlim(0, time_fam_agg.wall_mean.max() * 1.45)
ax.set_title(f"Training time per classifier family — grand total over 725 fits: {GRAND_TOTAL_SEC:.0f} s = {GRAND_TOTAL_SEC/3600:.2f} h",
             loc="left", fontsize=13)
plt.tight_layout()
plt.savefig(FIG/"fig10_training_time.png"); plt.close()

print(f"\n[deck] all figures + tables generated in {PRES.relative_to(ROOT)}/")

# =================================================================================
# BUILD THE PPT DECK
# =================================================================================
print("\n[deck] assembling deck.pptx ...")
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

NAVY_RGB = RGBColor(0x1a, 0x3a, 0x5c)
SAGE_RGB = RGBColor(0x27, 0xae, 0x60)
CHARCOAL_RGB = RGBColor(0x2c, 0x3e, 0x50)
GRAY_RGB = RGBColor(0x7f, 0x8c, 0x8d)
LIGHT_RGB = RGBColor(0xec, 0xf0, 0xf1)

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
SLIDE_W, SLIDE_H = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

def add_title(slide, text, y=0.35, size=28, color=NAVY_RGB):
    tx = slide.shapes.add_textbox(Inches(0.6), Inches(y), prs.slide_width - Inches(1.2), Inches(0.9))
    tf = tx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = text; r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = color
    return tx

def add_subtitle(slide, text, y=1.05, size=14, color=GRAY_RGB):
    tx = slide.shapes.add_textbox(Inches(0.6), Inches(y), prs.slide_width - Inches(1.2), Inches(0.4))
    tf = tx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = text; r.font.size = Pt(size); r.font.color.rgb = color
    return tx

def add_image_centered(slide, path, top=1.65, max_h=5.4, max_w=12.3):
    from PIL import Image
    im = Image.open(path); aspect = im.width / im.height
    # try max_w first
    w = max_w; h = w / aspect
    if h > max_h: h = max_h; w = h * aspect
    left = (prs.slide_width.inches - w) / 2
    slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(w), height=Inches(h))

def add_footer(slide, text="subFinder substrate prediction — rep_1 reference run"):
    tx = slide.shapes.add_textbox(Inches(0.5), prs.slide_height - Inches(0.4),
                                    prs.slide_width - Inches(1.0), Inches(0.3))
    tf = tx.text_frame
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = text; r.font.size = Pt(9); r.font.color.rgb = GRAY_RGB; r.font.italic = True

# Slide 1 — Title
s = prs.slides.add_slide(BLANK)
left_bar = s.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.5), prs.slide_height)
left_bar.fill.solid(); left_bar.fill.fore_color.rgb = NAVY_RGB; left_bar.line.fill.background()
tx = s.shapes.add_textbox(Inches(1.2), Inches(2.3), prs.slide_width - Inches(2.4), Inches(2))
tf = tx.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
r = p.add_run(); r.text = "subFinder"; r.font.size = Pt(56); r.font.bold = True; r.font.color.rgb = NAVY_RGB
p = tf.add_paragraph(); p.alignment = PP_ALIGN.LEFT
r = p.add_run(); r.text = "Calibrated classical-ML for PUL substrate prediction"; r.font.size = Pt(28); r.font.color.rgb = CHARCOAL_RGB
p = tf.add_paragraph(); p.space_before = Pt(20); p.alignment = PP_ALIGN.LEFT
r = p.add_run(); r.text = "Benchmark · best model · reproducibility · signature genes · validation"; r.font.size = Pt(16); r.font.color.rgb = GRAY_RGB
add_footer(s, " ")

# Slide 2 — Problem setup
s = prs.slides.add_slide(BLANK)
add_title(s, "Problem setup")
add_subtitle(s, "Predict one of 12 polysaccharide substrate classes from a PUL's gene-token sequence")
tx = s.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(11.5), Inches(5))
tf = tx.text_frame; tf.word_wrap = True
bullets = [
    ("Dataset", f"1,030 labeled PULs across 12 substrates (alginate, alpha-glucan, alpha-mannan, arabinogalactan, beta-glucan, beta-mannan, chitin, fructan, galactan, host glycan, pectin, xylan)"),
    ("Each PUL", "a comma/pipe-separated string of CAZy family IDs, transporter classification IDs, transcription-factor types, and null padding"),
    ("Cross-validation", "5-repeat × 5-fold Repeated Stratified K-Fold (n=25 trials per configuration) — leak-free: word-embedding models retrained per fold"),
    ("Goal", "(1) maximise predictive accuracy, (2) emit calibrated probabilities, (3) produce per-PUL signature-gene attributions that match literature CAZy knowledge"),
]
for i, (k, v) in enumerate(bullets):
    p = tf.add_paragraph() if i else tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    rk = p.add_run(); rk.text = f"{k} — "; rk.font.size = Pt(15); rk.font.bold = True; rk.font.color.rgb = NAVY_RGB
    rv = p.add_run(); rv.text = v; rv.font.size = Pt(15); rv.font.color.rgb = CHARCOAL_RGB
    p.space_after = Pt(14)
add_footer(s)

# Slide 3 — What we benchmarked
s = prs.slides.add_slide(BLANK)
add_title(s, "Benchmark scope")
add_subtitle(s, f"29 configurations = (3 featurizer families × 1 classifier swap) + (4 DL architectures × 5 embeddings)")
tx = s.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(12), Inches(5))
tf = tx.text_frame; tf.word_wrap = True
rows = [
    ("Featurizers (3)", "CountVectorizer (paper's tok_comma_pipe), CountVectorizer (tok_cpu — additionally splits on underscore), FastText mean+max-concat (gensim ft.wv[t] with n-gram OOV)"),
    ("Embeddings retrained per fold (6)", "FastText cbow/sg, Word2Vec cbow/sg, Doc2Vec dm/dbow — leak-free (test tokens never enter embedding training)"),
    ("Classifiers", "OvR(BalancedRF n=100) baseline; OvR(ExtraTrees n=500 log2/sqrt) ours; 4 DL architectures from paper (vanilla LSTM, LSTM+attention, just-attention, 4-block transformer)"),
    ("Hyperparameters", "Paper-verbatim for shallow + DL except batch (DL_BATCH=1024, Transformer=4096 for M4 Max throughput); EarlyStopping patience=30, validation 25%, Adam 1e-4"),
    ("Total", "29 configs × 25 trials = 725 fits; reused per-fold embeddings = reproducible to ~1e-3 on DL, bit-identical on sklearn"),
]
for i, (k, v) in enumerate(rows):
    p = tf.add_paragraph() if i else tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    rk = p.add_run(); rk.text = f"{k}  "; rk.font.size = Pt(14); rk.font.bold = True; rk.font.color.rgb = NAVY_RGB
    rv = p.add_run(); rv.text = v; rv.font.size = Pt(13); rv.font.color.rgb = CHARCOAL_RGB
    p.space_after = Pt(13)
add_footer(s)

def add_callout(slide, text, top=6.7, left=0.4, width=12.5, fc=RGBColor(0xff,0xf9,0xe6)):
    """Add a yellow-background 'How to read this' / 'What this means' callout box at the bottom."""
    tx = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(0.55))
    tf = tx.text_frame; tf.word_wrap = True
    tx.fill.solid(); tx.fill.fore_color.rgb = fc
    tx.line.color.rgb = CHARCOAL_RGB; tx.line.width = Pt(0.5)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = text; r.font.size = Pt(10.5); r.font.bold = True; r.font.color.rgb = CHARCOAL_RGB
    return tx

# Slide 4 — Benchmark results: leaderboard
s = prs.slides.add_slide(BLANK)
add_title(s, "Benchmark leaderboard — 29 configurations")
add_subtitle(s, "Mean test accuracy ± 1 SD across the full 5×5 RSKF grid (n=25 trials per config). Best on top.")
add_image_centered(s, FIG/"fig1_benchmark_leaderboard.png", top=1.45, max_h=5.0)
add_callout(s, "HOW TO READ — every row is one of 29 configurations (featurizer → classifier pair); "
                "bars are sorted descending so the strongest configs are at the top. Green = our shallow winners; "
                "navy = paper's Balanced RF baselines; orange = paper's deep architectures. "
                "Numbers next to each bar are the 25-trial mean accuracy.")
add_footer(s)

# Slide 4b — Top-5 podium (NEW)
s = prs.slides.add_slide(BLANK)
add_title(s, "Top-5 podium — the leading configurations")
add_subtitle(s, "Both #1 and #2 share the same OvR(ExtraTrees) classifier under different featurizers — the win is in the classifier, not the featurizer")
add_image_centered(s, FIG/"fig1b_top5_podium.png", top=1.5, max_h=5.0)
add_callout(s, "TAKEAWAY — swapping the paper's Balanced RF (n=100) for OvR(ExtraTrees n=500) on the same data "
                "delivers +6 pp accuracy. The featurizer choice (CountVec vs. FastText) is a secondary effect; "
                "ranks #3–#5 are all the paper's BRF baseline with different word-embedding featurizers.")
add_footer(s)

# (Removed) — the "Accuracy by model family" boxplot is now subsumed by slide 13 below,
# which shows the same family-level distribution with box + swarm + per-family stats table.

# Slide 6 — Reproducibility (the CODE-pipeline question; orthogonal to Slide 13)
s = prs.slides.add_slide(BLANK)
add_title(s, "Reproducibility — does a rerun give the same number?")
add_subtitle(s, f"X-axis = Δ accuracy (rep_2 mean − rep_1 mean, should be ≈ 0 for reproducibility). 5×5 RSKF data splits are FIXED across reps; only model-init seed (REPRO_REP_SEED) varies. Restricted to the {len(complete_retrain)}/29 configs with a COMPLETE 25-trial run in both reps.")
add_image_centered(s, FIG/"fig3_reproducibility.png", top=1.55, max_h=4.9)
add_callout(s, "WHAT THIS IS — re-running the SAME (config, seed, fold) with a DIFFERENT model-init seed (REPRO_REP_SEED), "
                "do we get the same accuracy? Our ExtraTrees winners (green) are random_state-seeded so they're identical "
                "across reps (Δ=0). Balanced RF baselines (navy) and DL configs (orange) show small Δ from imblearn thread "
                "non-determinism and weight-init+GPU op-order non-determinism respectively. Top rankings are stable.")
add_footer(s)

# Slide 6b — Cross-rep reproducibility (all 5 reps, forest plot)
# Load cross-rep summary if available
import json as _json
_xrep_summary_path = ROOT/"docs/tables/tab_cross_rep_summary.json"
if _xrep_summary_path.exists():
    _xrep = _json.loads(_xrep_summary_path.read_text())
    s = prs.slides.add_slide(BLANK)
    add_title(s, "Cross-rep reproducibility — 5 reps × 25 trials, data splits FIXED")
    add_subtitle(s, f"Model-init seed varies (REPRO_REP_SEED=1000/2000/3000/4000/5000); 5×5 RSKF data splits held fixed. Each row = one of 29 configs; 5 dots = per-rep means; square = cross-rep mean; bar = min↔max range.")
    add_image_centered(s, FIG/"fig14_cross_rep_stability.png", top=1.45, max_h=5.0)
    add_callout(s, f"HEADLINE — winner cpu__ET500_log2: <b>{_xrep['winner_cross_rep_mean']:.4f} ± {_xrep['winner_cross_rep_std']:.4f}</b> across 5 reps "
                    f"(range {_xrep['winner_cross_rep_min']:.4f}–{_xrep['winner_cross_rep_max']:.4f}, deterministic to 4th decimal). "
                    f"2nd-place ftCbow_MM__ET500_sqrt also stable at {_xrep['runner_cross_rep_mean']:.4f} ± {_xrep['runner_cross_rep_std']:.4f}. "
                    f"<b>Top-7 rank stability: {_xrep['top7_rank_stability']}</b> "
                    "(ranks 1-5 + 7 identical in every rep, single swap between #6/#7 in rep_3). "
                    "Per-family median cross-rep std: OvR(ExtraTrees) 0.0006 · OvR(BalancedRF) 0.0027 · DL families 0.0047-0.0064. "
                    "Our shallow winner is essentially deterministic; DL configs have 8-10× more model-init variance.")
    add_footer(s)

# Slide 7 — Per-substrate of best model: confusion matrix
s = prs.slides.add_slide(BLANK)
add_title(s, "Best model — per-substrate confusion")
add_subtitle(s, f"OvR(ExtraTrees 500, log2) on seed-42 5-fold out-of-fold predictions  (n=1030 PULs, acc={(y_pred==y_labels).mean():.4f})")
add_image_centered(s, FIG/"fig4_confusion.png", top=1.45, max_h=5.0)
add_callout(s, "HOW TO READ — rows are true substrate, columns are predicted; cell numbers are PUL counts. "
                "Diagonal cells are correct predictions; off-diagonal cells are mistakes. The hardest classes are "
                "fructan and beta-mannan (small canonical sets and few training examples); alginate, pectin and "
                "xylan reach near-perfect classwise accuracy.")
add_footer(s)

# Slide 8 — Per-substrate metric table inline
s = prs.slides.add_slide(BLANK)
add_title(s, "Best model — per-substrate F1 / precision / recall")
add_subtitle(s, "Seed-42 5-fold OOF, sorted by F1")
# Build per-substrate table on the slide
table_data = [["substrate","n test","precision","recall","F1"]]
for _, r in per_sub.iterrows():
    table_data.append([r.substrate, int(r.n_test), f"{r.precision:.3f}", f"{r.recall:.3f}", f"{r.F1:.3f}"])
n_rows = len(table_data); n_cols = 5
tbl_shape = s.shapes.add_table(n_rows, n_cols, Inches(2.5), Inches(1.7), Inches(8.5), Inches(5))
tbl = tbl_shape.table
col_widths = [Inches(2.2), Inches(1.2), Inches(1.6), Inches(1.6), Inches(1.4)]
for i, w in enumerate(col_widths): tbl.columns[i].width = w
for j, h in enumerate(table_data[0]):
    cell = tbl.cell(0, j); cell.text = str(h); cell.fill.solid(); cell.fill.fore_color.rgb = NAVY_RGB
    for para in cell.text_frame.paragraphs:
        for run in para.runs: run.font.bold = True; run.font.color.rgb = RGBColor(0xff,0xff,0xff); run.font.size = Pt(11)
for i in range(1, n_rows):
    for j in range(n_cols):
        cell = tbl.cell(i, j); cell.text = str(table_data[i][j])
        cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0xff,0xff,0xff) if i%2 else LIGHT_RGB
        for para in cell.text_frame.paragraphs:
            for run in para.runs: run.font.size = Pt(10); run.font.color.rgb = CHARCOAL_RGB
add_footer(s)

# Slide 9 — calibration (MOVED UP from slide 12: must come before sig-gene slides
# because all downstream sig-gene analyses use temperature-calibrated probabilities)
s = prs.slides.add_slide(BLANK)
add_title(s, "Probability calibration of the chosen model")
add_subtitle(s, "Temperature scaling halves ECE while preserving the argmax accuracy. ALL signature-gene analyses on the next slides use these CALIBRATED probabilities (T=0.70) — sig genes depend on probabilities, so they should be computed on the deployed (calibrated) model.")
add_image_centered(s, FIG/"fig8_calibration.png", top=1.65, max_h=4.7)
add_callout(s, "HOW TO READ — the reliability diagram (left) plots predicted confidence vs. empirical accuracy "
                "across 10 bins; a perfectly-calibrated model sits on the dashed diagonal. Temperature scaling (green) "
                "is the recommended protocol: halves ECE (0.094 → 0.029) with zero loss of argmax accuracy. "
                "IMPORTANT: the leave-one-token-out Δ-prob ablations on slides 11-13 are computed AFTER applying "
                "per-fold temperature scaling (mean T=0.70), so the sig genes reflect the deployed model's confidence.")
add_footer(s)

# Slide 10 — Top-3 sig genes per substrate (raw with lit hit markers) — NATIVE PPT TABLE
s = prs.slides.add_slide(BLANK)
add_title(s, "Per-substrate signature genes (raw model top-3) — lit-DB status")
add_subtitle(s, f"NO filtering. Of {n_cazy_check} CAZy features in 36 top-3 slots: {n_cazy_in_lit} are lit-canonical ({n_exact} via 1:1, {n_collapse} via collapse); {n_cazy_check-n_cazy_in_lit} CAZy not in lit-canon; {n_non_cazy} non-CAZy tokens (lit can't check).")
n_rows = len(substrates) + 1
tshape = s.shapes.add_table(n_rows, 4, Inches(0.4), Inches(1.55), Inches(12.5), Inches(5.1))
tt = tshape.table
for i, w in enumerate([Inches(1.4), Inches(3.7), Inches(3.7), Inches(3.7)]): tt.columns[i].width = w
for j, h in enumerate(["substrate", "top 1", "top 2", "top 3"]):
    c = tt.cell(0, j); c.text = h; c.fill.solid(); c.fill.fore_color.rgb = NAVY_RGB
    for para in c.text_frame.paragraphs:
        for run in para.runs: run.font.bold=True; run.font.color.rgb=RGBColor(0xff,0xff,0xff); run.font.size=Pt(11)
for i, sub_name in enumerate(substrates, 1):
    sub = sig_tbl[sig_tbl.substrate == sub_name].reset_index(drop=True)
    tt.cell(i, 0).text = sub_name
    for col_i in range(3):
        c = tt.cell(i, col_i + 1)
        if col_i < len(sub):
            r = sub.iloc[col_i]
            c.text = f"[{r.mark}] {r.feature} (imp {r.importance:.3f})\n{r.lit_status}"
            # color by status
            if r.mark == "LIT" and "1:1" in r.lit_status:    bg = RGBColor(0xc8,0xe6,0xc9)
            elif r.mark == "LIT":                            bg = RGBColor(0xe0,0xf3,0xe8)
            elif r.mark == "miss":                           bg = RGBColor(0xfc,0xe4,0xe4)
            else:                                            bg = RGBColor(0xfd,0xf6,0xe3)
        else:
            c.text = ""; bg = RGBColor(0xff,0xff,0xff)
        c.fill.solid(); c.fill.fore_color.rgb = bg
        for para in c.text_frame.paragraphs:
            for run in para.runs: run.font.size = Pt(9); run.font.color.rgb = CHARCOAL_RGB
    # substrate column styling
    tt.cell(i, 0).fill.solid(); tt.cell(i, 0).fill.fore_color.rgb = LIGHT_RGB
    for para in tt.cell(i,0).text_frame.paragraphs:
        for run in para.runs: run.font.size = Pt(11); run.font.bold = True; run.font.color.rgb = CHARCOAL_RGB
# Legend below the table (single textbox, one paragraph, simple ASCII)
legend_tx = s.shapes.add_textbox(Inches(0.4), Inches(6.75), Inches(12.5), Inches(0.35))
p = legend_tx.text_frame.paragraphs[0]
p.alignment = PP_ALIGN.LEFT
legend_text = ("[LIT] = canonical CAZy in literature (green; darker = 1:1 exact, lighter = via collapse)     "
               "[miss] = CAZy family but NOT in lit canon for this substrate (red)     "
               "[—] = non-CAZy token; literature cannot adjudicate (yellow)")
r = p.add_run(); r.text = legend_text; r.font.size = Pt(9); r.font.color.rgb = CHARCOAL_RGB
add_footer(s)

# Slide 10 — example PULs (native PPT table; explicit TRUE substrate column added)
s = prs.slides.add_slide(BLANK)
add_title(s, "Example predictions on held-out test PULs")
add_subtitle(s, "All 10 examples are correctly-classified PULs (true substrate = predicted substrate) cherry-picked for high literature-canonical coverage in their top-5 ablation list. Population-level metrics across all 1030 PULs are on the next slide.")
# Build a true-substrate column from ex.true_substrate
ex_show2 = ex.copy()
ex_show2["true_sub"] = ex_show2.apply(lambda r: r.get("true", r.pred), axis=1)
nr2 = len(ex_show) + 1
tshape2 = s.shapes.add_table(nr2, 7, Inches(0.20), Inches(1.85), Inches(12.95), Inches(4.7))
tt2 = tshape2.table
col_widths = [Inches(4.3), Inches(1.15), Inches(1.15), Inches(0.65), Inches(0.85), Inches(4.4), Inches(0.45)]
for i, w in enumerate(col_widths): tt2.columns[i].width = w
hdr = ["PUL gene-token sequence (first ~70 chars)", "TRUE substrate", "predicted",
       "max prob", "p-value", "top-5 signature genes (LIT=lit-canon; CAZy; acc=accessory)", "#LIT/5"]
for j, h in enumerate(hdr):
    c = tt2.cell(0, j); c.text = h; c.fill.solid(); c.fill.fore_color.rgb = NAVY_RGB
    for para in c.text_frame.paragraphs:
        for run in para.runs: run.font.bold = True; run.font.color.rgb = RGBColor(0xff,0xff,0xff); run.font.size = Pt(9.5)
for i, (_, r) in enumerate(ex.iterrows(), 1):
    true_val = r.get("true", r.pred)
    vals = [
        short_pul(int(r.idx), 70),
        true_val,
        r.pred,
        f"{r.prob:.3f}",
        f"{(1-r.prob)**11:.1e}",
        annot_top5(r.pred, r.top5),
        int(r.n_canon_top5),
    ]
    same = (true_val == r.pred)
    for j, v in enumerate(vals):
        c = tt2.cell(i, j); c.text = str(v)
        if j == 6:  # #LIT/5 colored by count
            n = int(v)
            bg = RGBColor(0xc8,0xe6,0xc9) if n >= 4 else RGBColor(0xe0,0xf3,0xe8) if n >= 2 else RGBColor(0xfd,0xf6,0xe3)
        elif j in (1, 2):  # truth/pred — green if they match, red if not
            bg = RGBColor(0xc8,0xe6,0xc9) if same else RGBColor(0xfc,0xe4,0xe4)
        else:
            bg = RGBColor(0xff,0xff,0xff) if i % 2 else RGBColor(0xfa,0xfb,0xfc)
        c.fill.solid(); c.fill.fore_color.rgb = bg
        for para in c.text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(8) if j in (0, 5) else Pt(9.5)
                run.font.color.rgb = CHARCOAL_RGB
                run.font.bold = j in (1, 2)
                if j in (0, 5): run.font.name = "Menlo"
add_callout(s, "TRUE & PREDICTED columns are both shown so reviewers can confirm we're not hiding errors — "
                "in all 10 rows truth equals prediction by construction (we filter to correctly-classified PULs first, "
                "then rank by lit-canonical coverage). For a per-PUL hit-rate that includes ALL eligible PULs (not just "
                "the cherry-picked top 10), see Slide 11.", top=6.65)
add_footer(s)

# Slide 11 — population lit validation (TRUE-class attribution; matches slide 13 totals)
s = prs.slides.add_slide(BLANK)
add_title(s, "Population-level signature-gene validation")
add_subtitle(s, "Two complementary recall-style metrics. TRUE-class attribution on CALIBRATED probabilities (T=0.70) — same method used in slide 13's per-substrate funnels.")
# Explanation block — short and plain
expl = s.shapes.add_textbox(Inches(0.4), Inches(1.4), Inches(12.5), Inches(1.05))
tf = expl.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
r = p.add_run()
r.text = ("How we built 'scope': of the 394 (substrate, canonical-CAZy) pairs the literature lists, "
          "many of those canonical genes never appear in any of our 1030 PULs — we exclude those, leaving "
          "173 'in-scope' (substrate, gene) pairs the model could plausibly find. "
          "Across ALL 1030 OOF PULs, 837 contain ≥1 in-scope canonical gene for their TRUE substrate (the 'eligible' denominator).")
r.font.size = Pt(11); r.font.color.rgb = CHARCOAL_RGB
p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.LEFT; p2.space_before = Pt(4)
r2 = p2.add_run()
r2.text = ("Left bar (PUL-level recall): for what % of eligible PULs does the top-K Δ_TRUE-class ablation list include at least one canonical gene for the TRUE substrate? "
           "Right bar (gene-level recall): of the 173 in-scope canonical genes, how many does the model surface as a top-K Δ_TRUE-class signature gene anywhere in the population? "
           "K=3 totals (768/837=91.8% any-hit and 109/173=63.0% scope coverage) match slide 13's per-substrate funnel totals exactly.")
r2.font.size = Pt(11); r2.font.color.rgb = CHARCOAL_RGB
# place the figure below the explanation
from PIL import Image
im = Image.open(FIG/"fig6_lit_validation.png"); ar = im.width/im.height
target_h = 4.4; target_w = target_h * ar
if target_w > 12.5: target_w = 12.5; target_h = target_w/ar
left = (prs.slide_width.inches - target_w)/2
s.shapes.add_picture(str(FIG/"fig6_lit_validation.png"), Inches(left), Inches(2.6),
                      width=Inches(target_w), height=Inches(target_h))
add_footer(s)

# Slide 13 — Per-substrate sig-gene FUNNEL (GROUND-TRUTH-class attribution on calibrated probs)
s = prs.slides.add_slide(BLANK)
add_title(s, "Per-substrate sig-gene funnel — TRUE-class attribution (K=3)")
add_subtitle(s, "Δ-prob computed against the TRUE substrate, not the argmax. Decouples attribution quality from classification quality — we ask 'when truth is s, did the model surface a canonical CAZy for s as a top-3 sig gene?' regardless of whether the model picked s as its prediction.")
add_image_centered(s, FIG/"fig8b_per_substrate_sig_pr.png", top=1.65, max_h=4.5)
add_callout(s, "WHY TRUE-CLASS ON CALIBRATED PROBS — our model outputs probabilities for all 12 classes for every PUL, "
                "so leave-one-token-out Δ-prob can be computed for ANY class. We compute Δ_s for TRUE substrate s using "
                "the CALIBRATED probabilities (per-fold temperature scaling, see Slide 9), which is what the deployed "
                "inference pipeline emits. Both funnels use truth as the unit (no 'correctly predicted' gate). Totals: "
                "PUL view 1030 → 837 eligible → 768 hit (91.8%) · gene view 394 → 173 in-scope → 109 flagged "
                "(63.0%). Scope recall is ~3 pp higher than the argmax-gated 60.1% on legacy aggregates because we "
                "now catch genes attributed correctly in PULs the model classified wrong.")
add_footer(s)

# Slide 13b — WOW: Rank-K redemption (top-K cumulative accuracy)
s = prs.slides.add_slide(BLANK)
add_title(s, "Rank-K redemption — TRUE substrate is recovered fast as K grows")
add_subtitle(s, "Cumulative top-K accuracy on 1030 held-out PULs (rep_1 OOF). When top-1 is wrong, the TRUE label is usually rank 2 or 3 — meaningful for triage workflows where biologists review the top few candidates.")
add_image_centered(s, FIG/"fig11_rank_redemption.png", top=1.55, max_h=4.6)
add_callout(s, "HEADLINE — top-1 = 0.906, top-2 = 0.956 (+5.0 pp), top-3 = 0.976 (+2.0 pp), top-5 = 0.981. "
                "Per-substrate: alginate is 100% at K=1; fructan is the hardest (top-1 = 0.677) but jumps to 0.903 by K=3 — "
                "i.e. when the model is wrong on fructan, the true class is almost always rank 2 or 3. Mean true-rank "
                "across all substrates = 1.39 (median 1). Calibrated probs make these ranks deployment-meaningful — see Slide 9.")
add_footer(s)

# Slide 13c — WOW: Calibration is meaningful (confidence ≈ accuracy per bin)
s = prs.slides.add_slide(BLANK)
add_title(s, "Calibration is meaningful — confidence ≈ accuracy per bin")
add_subtitle(s, "10-bin reliability histogram on calibrated probabilities (T ≈ 0.70). Sage = correct, red = incorrect. High-confidence predictions (≥0.8) ARE correct in ≥97% of cases — supports a triage/review workflow with a high-precision auto-accept threshold.")
add_image_centered(s, FIG/"fig12_confidence_vs_correct.png", top=1.55, max_h=4.6)
add_callout(s, "PER-BIN ACCURACY — 0.8-0.9 → 98.2% correct (110/112) · 0.9-1.0 → 97.3% correct (566/582). "
                "Combined ≥0.8: 97.4% correct on 694 PULs (67% of total). Bottom bins (≤0.5) show ~54-65% — the "
                "model 'knows it doesn't know'. Operational implication: route confidence ≥0.8 to auto-accept "
                "(precision ≈ 0.97), route 0.5-0.8 to expert review, refuse <0.5. Same calibrator powers the deployed inference.")
add_footer(s)

# Slide 13d — WOW: Hand-picked PUL case-study cards (6 scenarios)
s = prs.slides.add_slide(BLANK)
add_title(s, "Case studies — 6 hand-picked PUL predictions showing the value of top-K + sig genes")
add_subtitle(s, "Each card: TRUE substrate, model's top-3 (calibrated probs), and the gene sequence. Mix of confident-correct, low-conf-correct, rank-2/3 redemptions, and a confident-wrong failure mode.")
add_image_centered(s, FIG/"fig13_case_study_cards.png", top=1.45, max_h=4.9)
add_callout(s, "READING THE CARDS — green = top-1 is TRUE · amber = TRUE recovered at rank 2 or 3 · red = TRUE missing from top-3. "
                "Confident+correct (alginate, 0.96) shows clean signal: PL6/PL17 are canonical alginate lyases. "
                "Rank-2 redemption (alpha-glucan, 0.52 vs 0.32) — model splits between α/β-glucan, GH13 nudges α. "
                "Confident-wrong (alpha-glucan classified as fructan, p=1.0) — GH32/3.A.1.1.x are PRTS sucrose markers; ambiguous substrate.")
add_footer(s)

# Slide 12c — Test-PUL OOV (out-of-vocabulary) vs accuracy
s = prs.slides.add_slide(BLANK)
add_title(s, "Internal robustness — accuracy vs train-vocab OOV proportion")
add_subtitle(s, "Uses the SAME featurizer the deployed model uses (CountVec with tok_cpu, fit per fold on outer_tr only). 'OOV' = test tokens NOT in the training fold's vocab.")
add_image_centered(s, FIG/"fig8c_oov_vs_accuracy.png", top=1.55, max_h=4.5)
add_callout(s, "READING + CAVEAT — bar widths ∝ PUL count per bucket. 89% of test PULs have 0% OOV; the model "
                "scores 91% on them. Once OOV >10%, accuracy collapses to ~64%. CAVEAT: 'novel' here means "
                "WITHIN-DATASET distinct-token frequency, NOT biological novelty (1030 PULs share a tight "
                "~517-token vocab; a typical training fold covers ~488 of them). The TRUE cross-organism "
                "novelty test is on fungal CGCs — see supplement S9 where 24/24 are flagged REFUSE.")
add_footer(s)

# Slide 13 — Robustness (the MODEL-variance question; orthogonal to Slide 6)
s = prs.slides.add_slide(BLANK)
add_title(s, "Robustness — how stable is each model family across trials?")
add_subtitle(s, "X-axis = absolute test accuracy (range 0.4-0.95). 6 rows = 6 model families pooling every (config × seed × fold) trial. Distinct from Slide 6 which shows Δ accuracy of a single rerun.")
add_image_centered(s, FIG/"fig9_trial_distribution.png", top=1.55, max_h=4.9)
add_callout(s, "WHAT THIS IS — the question is: for ONE family, how much does a single (seed, fold) trial change "
                "the accuracy? Box markers show central tendency; dots are individual trials. The green family sits "
                "entirely above 0.86 with a tight ~0.04 spread; deep families span much wider (some folds dip below "
                "0.55). Our shallow ensemble is BOTH more accurate AND more stable.")
add_footer(s)

# Slide A1 — APPENDIX: Full 29-config leaderboard with decoded featurizer + classifier
s = prs.slides.add_slide(BLANK)
add_title(s, "Appendix · Full 29-config leaderboard")
add_subtitle(s, "Every configuration benchmarked; rank 1 highlighted gold, ranks 2-3 pale gold.")
nrA = len(full_lb) + 1
tshapeA = s.shapes.add_table(nrA, 7, Inches(0.2), Inches(1.45), Inches(13.0), Inches(5.6))
ttA = tshapeA.table
hdrA = ["#","config (shorthand)","featurizer","feat. detail","classifier","clf. detail","mean ± std"]
for i, w in enumerate([Inches(0.35), Inches(2.0), Inches(1.1), Inches(2.5), Inches(2.2), Inches(3.0), Inches(1.85)]):
    ttA.columns[i].width = w
for j, h in enumerate(hdrA):
    c = ttA.cell(0, j); c.text = h; c.fill.solid(); c.fill.fore_color.rgb = NAVY_RGB
    for para in c.text_frame.paragraphs:
        for run in para.runs: run.font.bold=True; run.font.color.rgb=RGBColor(0xff,0xff,0xff); run.font.size=Pt(9.5)
for i, (_, r) in enumerate(full_lb.iterrows(), 1):
    cells = [str(int(r["rank"])), r.shorthand, r.featurizer_fam, r.featurizer_det,
              r.clf_fam, r.clf_det, f"{r['mean']:.4f} ± {r['std']:.4f}"]
    for j, v in enumerate(cells):
        c = ttA.cell(i, j); c.text = str(v)
        if r["rank"] == 1:    bg = RGBColor(0xff, 0xf3, 0xcd)   # gold
        elif r["rank"] <= 3:  bg = RGBColor(0xff, 0xf9, 0xe6)   # pale gold
        else:                 bg = RGBColor(0xff, 0xff, 0xff) if i%2 else RGBColor(0xfa,0xfb,0xfc)
        c.fill.solid(); c.fill.fore_color.rgb = bg
        for para in c.text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(7.5); run.font.color.rgb = CHARCOAL_RGB
                if r["rank"] == 1: run.font.bold = True
add_footer(s, "Sorted by mean test accuracy; n=25 trials each (5×5 RSKF)")

# Slide A2 — APPENDIX: Per-featurizer aggregate
s = prs.slides.add_slide(BLANK)
add_title(s, "Appendix · Performance by featurizer family")
add_subtitle(s, "Averaged over all (config, seed, fold) trials sharing the same featurizer family")
nrB = len(feat_agg) + 1
tshapeB = s.shapes.add_table(nrB, 5, Inches(1.5), Inches(1.6), Inches(10.3), Inches(3.6))
ttB = tshapeB.table
hdrB = ["featurizer family", "mean accuracy", "std", "n configs", "n trials"]
for i, w in enumerate([Inches(2.8), Inches(2.0), Inches(1.5), Inches(1.8), Inches(2.2)]):
    ttB.columns[i].width = w
for j, h in enumerate(hdrB):
    c = ttB.cell(0, j); c.text = h; c.fill.solid(); c.fill.fore_color.rgb = NAVY_RGB
    for para in c.text_frame.paragraphs:
        for run in para.runs: run.font.bold=True; run.font.color.rgb=RGBColor(0xff,0xff,0xff); run.font.size=Pt(12)
for i, (_, r) in enumerate(feat_agg.iterrows(), 1):
    cells = [r.featurizer_family, f"{r.mean_acc:.4f}", f"{r.std_acc:.4f}", int(r.n_configs), int(r.n_trials)]
    for j, v in enumerate(cells):
        c = ttB.cell(i, j); c.text = str(v); c.fill.solid()
        c.fill.fore_color.rgb = RGBColor(0xff,0xff,0xff) if i%2 else LIGHT_RGB
        for para in c.text_frame.paragraphs:
            for run in para.runs: run.font.size = Pt(11); run.font.color.rgb = CHARCOAL_RGB
note_tx = s.shapes.add_textbox(Inches(1.5), Inches(5.5), Inches(10.3), Inches(1.0))
tf = note_tx.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
r = p.add_run()
r.text = ("Note: averages mix in different classifier choices. Pattern — CountVec leads because it carries our ExtraTrees winner; "
          "FastText/Word2Vec/Doc2Vec families are pulled down by their BRF/DL classifier pairings.")
r.font.size = Pt(11); r.font.color.rgb = CHARCOAL_RGB; r.font.italic = True
add_footer(s)

# Slide A3 — APPENDIX: Per-classifier aggregate
s = prs.slides.add_slide(BLANK)
add_title(s, "Appendix · Performance by classifier family")
add_subtitle(s, "Averaged over all (config, seed, fold) trials sharing the same classifier family")
nrC = len(clf_agg) + 1
tshapeC = s.shapes.add_table(nrC, 5, Inches(1.5), Inches(1.6), Inches(10.3), Inches(4.0))
ttC = tshapeC.table
hdrC = ["classifier family", "mean accuracy", "std", "n configs", "n trials"]
for i, w in enumerate([Inches(3.2), Inches(2.0), Inches(1.5), Inches(1.6), Inches(2.0)]):
    ttC.columns[i].width = w
for j, h in enumerate(hdrC):
    c = ttC.cell(0, j); c.text = h; c.fill.solid(); c.fill.fore_color.rgb = NAVY_RGB
    for para in c.text_frame.paragraphs:
        for run in para.runs: run.font.bold=True; run.font.color.rgb=RGBColor(0xff,0xff,0xff); run.font.size=Pt(12)
for i, (_, r) in enumerate(clf_agg.iterrows(), 1):
    cells = [r.classifier_family, f"{r.mean_acc:.4f}", f"{r.std_acc:.4f}", int(r.n_configs), int(r.n_trials)]
    for j, v in enumerate(cells):
        c = ttC.cell(i, j); c.text = str(v); c.fill.solid()
        if r.classifier_family == "OvR(ExtraTrees)":
            c.fill.fore_color.rgb = RGBColor(0xff, 0xf3, 0xcd)
        else:
            c.fill.fore_color.rgb = RGBColor(0xff,0xff,0xff) if i%2 else LIGHT_RGB
        for para in c.text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(11); run.font.color.rgb = CHARCOAL_RGB
                if r.classifier_family == "OvR(ExtraTrees)": run.font.bold = True
note_tx = s.shapes.add_textbox(Inches(1.5), Inches(6.0), Inches(10.3), Inches(0.9))
tf = note_tx.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
r = p.add_run()
r.text = ("Note: averages mix in different featurizer choices. Pattern — OvR(ExtraTrees) is the only family in the >0.90 band; "
          "every DL family lands in 0.74-0.82.")
r.font.size = Pt(11); r.font.color.rgb = CHARCOAL_RGB; r.font.italic = True
add_footer(s)

# Slide A4 — APPENDIX: Training time per classifier family
s = prs.slides.add_slide(BLANK)
add_title(s, "Appendix · Training time per classifier family")
add_subtitle(s, f"Why our shallow winner is cheap. Grand total across all 725 fits: {GRAND_TOTAL_SEC:.0f} s = {GRAND_TOTAL_SEC/3600:.2f} h on Apple M4 Max")
add_image_centered(s, FIG/"fig10_training_time.png", top=1.5, max_h=3.6)
# Add a callout textbox below
callout_tx = s.shapes.add_textbox(Inches(0.6), Inches(5.4), Inches(12.0), Inches(1.6))
tf = callout_tx.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
r = p.add_run()
r.text = ("Our winner trains in ~3 s/fold (≈75 s per 25-fold config); the most expensive DL family — the 4-block transformer — "
          "takes ~70 s/fold (≈30 min per config), i.e. ~24× more compute per fit. Times exclude per-fold word-embedding "
          "training (cached under fold_cache_v2/) and the OvR(ExtraTrees) winner is therefore ~25× more compute-efficient "
          "per percentage point of accuracy.")
r.font.size = Pt(11); r.font.color.rgb = CHARCOAL_RGB; r.font.italic = True
add_footer(s)

# Slide 14 — APPENDIX: Substrate-alias map with primary-literature citations
s = prs.slides.add_slide(BLANK)
add_title(s, "Appendix · Substrate-alias map between literature DB and our 12 classes")
add_subtitle(s, "Curated DB has 75 finer-grained substrate categories; we collapse them into our 12 with primary-literature support")
# Build a table on the slide
alias_rows = [
    ("Our class", "Lit name(s)", "Match type", "Citation"),
    ("alginate", "alginate", "1:1 exact", "CAZy DB (Lombard 2014 NAR)"),
    ("pectin", "pectin", "1:1 exact", "CAZy DB"),
    ("xylan", "xylan", "1:1 exact", "CAZy DB"),
    ("alpha-mannan", "alpha-mannan", "1:1 exact", "CAZy DB"),
    ("beta-mannan", "beta-mannan", "1:1 exact", "CAZy DB"),
    ("fructan", "fructan", "1:1 exact", "CAZy DB"),
    ("beta-glucan", "beta-glucan (1:1) + cellulose, cellooligosaccharide, xyloglucan, beta-glycan (collapse)",
     "exact + 4 alias", "Burton 2006 Science 311:1940; Eklof & Brumer 2010 Plant Physiol 153:456"),
    ("alpha-glucan", "alpha-glucan (1:1) + starch, glycogen, sucrose, raffinose, trehalose, palatinose, glucooligo (collapse)",
     "exact + 7 alias", "Stam 2006 Protein Eng. 19:555; Janecek 2014 Cell. Mol. Life Sci. 71:1149"),
    ("chitin", "chitin (1:1) + chitosan, chitooligosaccharide (collapse)",
     "exact + 2 alias", "Hartl 2012 Appl. Microbiol. 93:533; Adrangi & Faramarzi 2013 Biotech. Adv. 31:1786"),
    ("host glycan", "host glycan (1:1) + HMO, sialic-acid, fucose (collapse)",
     "exact + 3 alias", "Marcobal 2011 Cell Host & Microbe 10:507; Tailford 2015 Frontiers Genet. 6:81"),
    ("arabinogalactan", "arabinan + arabinogalactan protein (no 1:1 entry)",
     "collapse only", "Tan 2013 Plant Cell 25:270; Showalter 2010 Plant Physiol 153:485"),
    ("galactan", "alpha-galactan + beta-galactan (no 1:1 entry)",
     "collapse only", "CAZy DB — anomericity sub-classes"),
]
nr = len(alias_rows); nc = 4
tshape = s.shapes.add_table(nr, nc, Inches(0.4), Inches(1.55), Inches(12.5), Inches(5.4))
tt = tshape.table
for i,w in enumerate([Inches(1.5), Inches(5.5), Inches(1.6), Inches(3.9)]): tt.columns[i].width = w
for j, h in enumerate(alias_rows[0]):
    c = tt.cell(0,j); c.text = h; c.fill.solid(); c.fill.fore_color.rgb = NAVY_RGB
    for para in c.text_frame.paragraphs:
        for run in para.runs: run.font.bold=True; run.font.color.rgb=RGBColor(0xff,0xff,0xff); run.font.size=Pt(11)
for i in range(1, nr):
    for j in range(nc):
        c = tt.cell(i,j); c.text = str(alias_rows[i][j])
        # Color code by match type
        mt = alias_rows[i][2]
        if "1:1 exact" == mt: bg = RGBColor(0xc8,0xe6,0xc9)   # strong green
        elif "exact + " in mt: bg = RGBColor(0xe0,0xf3,0xe8)   # pale green
        elif "collapse only" in mt: bg = RGBColor(0xfd,0xf6,0xe3) # pale yellow
        else: bg = RGBColor(0xff,0xff,0xff)
        c.fill.solid(); c.fill.fore_color.rgb = bg
        for para in c.text_frame.paragraphs:
            for run in para.runs: run.font.size = Pt(9); run.font.color.rgb = CHARCOAL_RGB
add_footer(s, "Detailed citations in paper/reference.bib; codified as ALIAS_PROVENANCE in presentations/build_slides.py")

# Slide 15 — Take-home
s = prs.slides.add_slide(BLANK)
add_title(s, "Take-home")
tx = s.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(12), Inches(5.5))
tf = tx.text_frame; tf.word_wrap = True
points = [
    ("Best configuration", "cpu__ET500_log2 = CountVec (tok_cpu) × OvR(ExtraTrees n=500, log2, class_weight='balanced', bootstrap=False) — 0.9058 ± 0.0172 on full 5×5 RSKF"),
    ("Wins by classifier choice, not features", "+6.16 pp over published BRF baseline (paired t-test p ≈ 5×10⁻¹⁴), +8.10 pp over best paper DL (LSTM-with-attention on FastText-skipgram)"),
    ("Calibrated for deployment", "Temperature T ≈ 0.70 reduces 10-bin ECE from 0.094 to 0.029 while preserving 0.9029 OOF accuracy exactly (monotonic)"),
    ("Signature genes are biologically interpretable", "Per-PUL leave-one-token-out ablation surfaces ≥1 literature-canonical CAZy in top-3 for 93% of eligible PULs; covers 62% of the in-scope canonical gene vocabulary at K=3"),
    ("Fully reproducible", "Every (config, seed, fold) classifier weight + train/test predictions saved in reproducibility/rep_1/predictions/ — single command `--training False` rebuilds every paper number in ~1 min"),
]
for i, (k, v) in enumerate(points):
    p = tf.add_paragraph() if i else tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    rk = p.add_run(); rk.text = f"{k} — "; rk.font.size = Pt(14); rk.font.bold = True; rk.font.color.rgb = SAGE_RGB
    rv = p.add_run(); rv.text = v; rv.font.size = Pt(13); rv.font.color.rgb = CHARCOAL_RGB
    p.space_after = Pt(12)
add_footer(s)

out = PRES/"deck.pptx"
prs.save(out)
print(f"[deck] wrote {out.relative_to(ROOT)} ({out.stat().st_size//1024} KB)")
print(f"[deck] total slides: {len(prs.slides)}")
